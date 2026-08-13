"""The main canvas: drop spectra straight onto it and see them.

Deliberately simple, replacing the earlier layout-first flow (define a figure
with N slots, then fill the slots). Here you just drag one or more datasets
from the browser onto the canvas and they are loaded and drawn. Arrangement
(overlay vs stacked) is a control you flip afterwards, not a decision you have
to make up front.

Loading happens on a worker thread: reading processed data off a network share
is slow enough to freeze the GUI if done inline. Each spectrum is drawn as it
arrives, so dropping several does not block on the slowest one.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path

import numpy as np
from matplotlib.backends.backend_qtagg import FigureCanvasQTAgg
from matplotlib.figure import Figure
from PySide6.QtCore import QObject, QRunnable, Qt, QThreadPool, Signal
from PySide6.QtWidgets import QFileDialog, QMenu, QVBoxLayout, QWidget

from ..domain.project import DEFAULT_PALETTE as PALETTE
from .dataset_model import MIME_DATASET

# Okabe-Ito is already the domain palette; reuse it so on-screen colours match
# whatever an eventual export produces.
# Yellow (#F0E442) is deliberately absent: it is nearly invisible as a thin
# line on white. Order puts high-contrast colours first so a two-spectrum
# comparison -- by far the common case -- is black against strong blue.
_FALLBACK_PALETTE = [
    "#000000", "#0072B2", "#D55E00", "#009E73",
    "#CC79A7", "#56B4E9", "#E69F00", "#8C564B",
]


# Real spectra in one figure can differ by many orders of magnitude (a 5 mM
# reference next to a 50 uM sample is ~1000x on its own), so the scale range
# has to span far more than a few decades. A 1000x ceiling silently clamped
# autoscale and made strong/weak pairs impossible to compare.
MIN_Y_SCALE = 1e-9
MAX_Y_SCALE = 1e9

LINE_STYLES = {
    "Solid": "-",
    "Dashed": "--",
    "Dotted": ":",
    "Dash-dot": "-.",
}


def _palette() -> list[str]:
    try:
        return list(PALETTE) or _FALLBACK_PALETTE
    except Exception:  # noqa: BLE001 - palette shape is not load-bearing here
        return _FALLBACK_PALETTE


class _PlotCanvas(FigureCanvasQTAgg):
    """Canvas subclass that guarantees mouse tracking.

    The wheel is handled through matplotlib's own ``scroll_event``, which the
    backend emits from FigureCanvasQT.wheelEvent. An earlier version tried to
    override that by assigning ``canvas.wheelEvent = fn`` on the INSTANCE --
    Qt dispatches virtuals through the class, so the assignment was silently
    ignored AND it shadowed the backend's own handler, killing scroll_event
    too. Hence: no wheel override here, just tracking so motion events arrive
    without a button held down.
    """

    def __init__(self, figure, owner):
        super().__init__(figure)
        self._owner = owner
        self.setMouseTracking(True)


@dataclass
class Trace:
    """One loaded 1D spectrum on the canvas.

    y_scale and y_offset are per-trace so a weak spectrum can be brought up to
    a strong one without touching the others -- adjusted with the scroll wheel
    over the plot, or typed exactly. line_width and color are per-trace too so
    a preference change can apply to all while still allowing overrides.
    """

    path: Path
    label: str
    ppm: np.ndarray
    intensity: np.ndarray
    color: str
    visible: bool = True
    y_scale: float = 1.0
    y_offset: float = 0.0
    line_width: float = 0.8
    line_style: str = "-"      # matplotlib style: '-', '--', ':', '-.'
    # 2D spectra carry a matrix plus both ppm axes. ppm/intensity stay empty
    # for these; is_2d is what the drawing code branches on.
    is_2d: bool = False
    matrix: "np.ndarray | None" = None
    ppm_f1: "np.ndarray | None" = None
    ppm_f2: "np.ndarray | None" = None
    contour_levels: int = 12
    # Label position in AXES fractions, so dragging a name moves it and it
    # still does not drift when the spectrum is rescaled.
    label_pos: "tuple[float, float] | None" = None
    pulse_program: str = ""
    nucleus: str = ""
    label_offset: "tuple[float, float]" = (0.0, 0.0)
    label_base_pos: "tuple[float, float] | None" = None
    is_difference: bool = False
    label_dx: float = 0.0     # extra label offset, axes fractions
    label_dy: float = 0.0
    pulse_program: str = ""
    nucleus: str = ""
    label_offset: "tuple[float, float]" = (0.0, 0.0)
    label_base_pos: "tuple[float, float] | None" = None
    is_difference: bool = False

    def display_label(self, show_pulprog: bool = True) -> str:
        """Name as drawn on the plot."""
        if show_pulprog and self.pulse_program:
            return f"{self.label}  ({self.pulse_program})"
        return self.label


class _LoadSignals(QObject):
    loaded = Signal(object, object, object, str)   # path, ppm, intensity, label
    loaded2d = Signal(object, object, str)         # path, payload dict, label
    failed = Signal(object, str)                   # path, message


class _LoadTask(QRunnable):
    """Reads one spectrum off the GUI thread.

    Pure I/O plus numpy: touches no Qt model state, so it is safe on a worker.
    The result is handed back by signal and applied on the GUI thread.
    """

    def __init__(self, reader, path: Path, label: str, dimensionality: int):
        super().__init__()
        self._reader = reader
        self._path = path
        self._label = label
        self._dim = dimensionality
        self.signals = _LoadSignals()

    def run(self) -> None:
        try:
            if self._dim == 2:
                spec = self._reader.read_2d(self._path)
                payload = {
                    "is_2d": True,
                    "matrix": np.asarray(spec.real, dtype=np.float64),
                    "ppm_f1": np.asarray(
                        spec.axis_f1.ppm_scale(), dtype=np.float64
                    ),
                    "ppm_f2": np.asarray(
                        spec.axis_f2.ppm_scale(), dtype=np.float64
                    ),
                }
                self.signals.loaded2d.emit(self._path, payload, self._label)
                return
            spec = self._reader.read_1d(self._path)
            ppm = np.asarray(spec.axis.ppm_scale(), dtype=np.float64)
            intensity = np.asarray(spec.real, dtype=np.float64)
            self.signals.loaded.emit(self._path, ppm, intensity, self._label)
        except Exception as exc:  # noqa: BLE001 - one bad file must not kill the drop
            self.signals.failed.emit(self._path, str(exc))


class SpectrumCanvas(QWidget):
    """A matplotlib canvas that accepts dataset drops and plots them.

    Public surface kept small on purpose: drop things on it, call
    set_arrangement / set_ppm_range / clear. No figure-layout model, no slots.
    """

    spectrumAdded = Signal(str)      # label
    loadFailed = Signal(str, str)    # path, message
    tracesChanged = Signal()
    cursorMoved = Signal(float, float)    # ppm, intensity under the cursor
    imageSaved = Signal(str)              # path written
    modeChanged = Signal(str)             # '1D' or '2D'
    dimensionalityRefused = Signal(str)   # message

    ARRANGEMENT_OVERLAY = "overlay"
    ARRANGEMENT_STACKED = "stacked"

    def __init__(self, reader=None, pool: QThreadPool | None = None, parent=None):
        super().__init__(parent)
        if reader is None:
            from ..infrastructure.nmrglue_reader import NmrglueReader

            reader = NmrglueReader()
        self._reader = reader
        self._pool = pool or QThreadPool.globalInstance()
        self._traces: list[Trace] = []
        self._arrangement = self.ARRANGEMENT_OVERLAY
        self._ppm_range: tuple[float, float] | None = None
        self._inflight: set = set()
        self._selected_index: int | None = None
        self._default_line_width = 0.8
        self._palette_override: list[str] | None = None
        # Appearance per SLOT (1st spectrum loaded, 2nd, ...), set from
        # Preferences. Kept here so a newly dropped spectrum picks up its
        # slot's colour/style/width immediately.
        from .preferences_dialog import default_styles
        self._slot_styles = default_styles()
        # Y limits are held FIXED once established. Without this matplotlib
        # autoscales, so multiplying a trace's intensity just rescales the
        # axis and the spectrum looks completely unchanged -- the reported
        # 'Y scale does nothing' bug.
        self._y_limits: tuple[float, float] | None = None
        self._show_grid = False
        # Grid spacing in ppm. None = let matplotlib choose.
        self._grid_spacing_ppm = None
        # Decimal places on the ppm axis labels. None = matplotlib default.
        self._x_decimals = None
        # Relative size of the on-plot spectrum names. Larger values also
        # space them further apart, so a long list does not overlap the
        # traces it is labelling.
        self._label_scale = 1.0
        self._show_pulprog = True
        self._crosshair = None
        self._crosshair_enabled = True
        self._drag_start = None
        self._label_artists = []
        self._label_drag = None
        self._pending_meta = {}
        self._f1_range = None
        self._f2_range = None
        self._last_mode = "1D"

        self._figure = Figure(figsize=(6, 4), tight_layout=False, facecolor="white")
        self._canvas = _PlotCanvas(self._figure, self)
        self._axes = self._figure.add_subplot(111)
        self._axes.set_facecolor("white")

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(self._canvas)

        self.setAcceptDrops(True)
        self.setContextMenuPolicy(Qt.CustomContextMenu)
        self.customContextMenuRequested.connect(self._on_context_menu)
        # Scroll wheel over the plot scales the SELECTED trace vertically.
        self._canvas.setFocusPolicy(Qt.WheelFocus)
        # matplotlib's own scroll event rather than overriding Qt's
        # wheelEvent: assigning an instance attribute over a C++ virtual
        # is not reliably dispatched by the Shiboken binding.
        self._canvas.mpl_connect("scroll_event", self._on_scroll)
        # Crosshair + drag-to-move use matplotlib's own event system rather
        # than Qt's, so the coordinates arrive already in DATA space.
        self._canvas.mpl_connect("motion_notify_event", self._on_mouse_move)
        self._canvas.mpl_connect("button_press_event", self._on_mouse_press)
        self._canvas.mpl_connect("button_release_event", self._on_mouse_release)
        self._canvas.mpl_connect("axes_leave_event", self._on_mouse_leave)
        self._redraw()

    # -- public state --------------------------------------------------------

    @property
    def traces(self) -> list[Trace]:
        return list(self._traces)

    def arrangement(self) -> str:
        return self._arrangement

    def set_arrangement(self, arrangement: str) -> None:
        if arrangement not in (self.ARRANGEMENT_OVERLAY, self.ARRANGEMENT_STACKED):
            return
        self._arrangement = arrangement
        self._y_limits = None   # stacking changes the frame entirely
        self._redraw()

    def set_ppm_range(self, left: float, right: float) -> None:
        """ppm axes descend, so left must be the HIGHER value."""
        if left <= right:
            return
        self._ppm_range = (left, right)
        self._redraw()

    def full_range(self) -> None:
        """Union of every loaded trace's ppm span (not the intersection) --
        show everything rather than only the overlap."""
        self._ppm_range = None
        self._redraw()

    def ppm_bounds(self) -> tuple[float, float] | None:
        """Union of every visible trace's ppm span, or None if nothing loaded.

        Union rather than intersection: the range control should be able to
        show everything that exists, not only the overlap.
        """
        visible = [t for t in self._traces if t.visible]
        if not visible:
            return None
        left = max(float(np.nanmax(t.ppm)) for t in visible)
        right = min(float(np.nanmin(t.ppm)) for t in visible)
        return (left, right)

    def clear(self) -> None:
        """Reset the canvas completely -- traces, selection, and ppm range.

        A true "start over", not just removing the lines: leaving a stale ppm
        range or selection behind after clearing is what makes the next drop
        appear on a nonsensical axis.
        """
        self._traces.clear()
        self._ppm_range = None
        self._f1_range = None
        self._f2_range = None
        self._selected_index = None
        self._y_limits = None
        self._redraw()
        self._emit_mode_if_changed()
        self.tracesChanged.emit()

    # -- selection and per-trace vertical scaling ----------------------------

    def selected_index(self) -> int | None:
        return self._selected_index

    def select_trace(self, index: int | None) -> None:
        """Select which trace the wheel / y-scale controls act on.

        None means 'nothing selected'; out-of-range indexes are ignored rather
        than raising, since a selection can outlive the trace it referred to.
        """
        if index is None:
            self._selected_index = None
        elif 0 <= index < len(self._traces):
            self._selected_index = index
        else:
            return
        self._redraw()
        self.tracesChanged.emit()

    def selected_trace(self) -> Trace | None:
        if self._selected_index is None:
            return None
        if not (0 <= self._selected_index < len(self._traces)):
            return None
        return self._traces[self._selected_index]

    def set_y_scale(self, index: int, scale: float) -> None:
        """Set one trace's vertical scale. Non-positive values are refused --
        a zero or negative scale would flatten or invert the spectrum, which
        is never what a scale control is meant to do."""
        if not (0 <= index < len(self._traces)):
            return
        if scale <= 0 or not np.isfinite(scale):
            return
        scale = max(MIN_Y_SCALE, min(float(scale), MAX_Y_SCALE))
        self._traces[index].y_scale = scale
        self._redraw()
        self.tracesChanged.emit()

    def nudge_y_scale(self, index: int, factor: float) -> None:
        """Multiply a trace's scale (wheel steps). Clamped to a sane range so
        a fast scroll cannot drive it to zero or to a value that overflows."""
        if not (0 <= index < len(self._traces)):
            return
        current = self._traces[index].y_scale
        new = current * float(factor)
        # Clamped to a range where the spectrum stays findable. The old
        # ceiling of 1e9 let a few scroll notches launch a trace far off
        # screen with no obvious way back.
        new = max(MIN_Y_SCALE, min(new, MAX_Y_SCALE))
        self.set_y_scale(index, new)

    def set_y_offset(self, index: int, offset: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        if not np.isfinite(offset):
            return
        self._traces[index].y_offset = float(offset)
        self._redraw()
        self.tracesChanged.emit()

    def nudge_y_offset(self, index: int, delta: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        self.set_y_offset(index, self._traces[index].y_offset + float(delta))

    def _on_scroll(self, event) -> None:
        """Wheel over the plot scales the selected trace vertically.

        With nothing selected the wheel does nothing rather than silently
        scaling an arbitrary trace -- guessing which spectrum the user meant
        would be worse than requiring a click first.
        """
        index = self._selected_index
        if index is None:
            return
        step = getattr(event, "step", 0) or 0
        if step == 0:
            button = getattr(event, "button", None)
            step = 1 if button == "up" else (-1 if button == "down" else 0)
        if step == 0:
            return
        # ~10% per notch, direction following the scroll.
        factor = 1.1 if step > 0 else (1.0 / 1.1)
        self.nudge_y_scale(index, factor)

    # -- appearance preferences ---------------------------------------------

    def set_default_line_width(self, width: float) -> None:
        """Applies to every trace, including ones already drawn."""
        if width <= 0 or not np.isfinite(width):
            return
        self._default_line_width = float(width)
        for trace in self._traces:
            trace.line_width = float(width)
        self._redraw()

    def default_line_width(self) -> float:
        return self._default_line_width

    def set_trace_line_width(self, index: int, width: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        if width <= 0 or not np.isfinite(width):
            return
        self._traces[index].line_width = float(width)
        self._redraw()

    def set_trace_color(self, index: int, color: str) -> None:
        if not (0 <= index < len(self._traces)):
            return
        if not color:
            return
        self._traces[index].color = color
        self._redraw()
        self.tracesChanged.emit()

    def set_palette(self, colors: list[str]) -> None:
        """Replace the colour cycle and recolour existing traces in order."""
        if not colors:
            return
        self._palette_override = list(colors)
        for i, trace in enumerate(self._traces):
            trace.color = colors[i % len(colors)]
        self._redraw()
        self.tracesChanged.emit()

    def _active_palette(self) -> list[str]:
        return self._palette_override or _palette()

    def set_trace_visible(self, index: int, visible: bool) -> None:
        if not (0 <= index < len(self._traces)):
            return
        self._traces[index].visible = bool(visible)
        # Deliberately does NOT refit the frame. Hiding a spectrum is a
        # viewing action, not a change to the data being framed: refitting
        # made everything jump in scale the moment a box was unticked, and
        # re-ticking it did not restore the previous view. Only adding or
        # removing spectra, or switching arrangement, re-fits.
        self._redraw()
        self.tracesChanged.emit()

    def remove_trace(self, path: Path) -> None:
        before = len(self._traces)
        self._traces = [t for t in self._traces if t.path != Path(path)]
        if len(self._traces) != before:
            self._y_limits = None   # re-fit to what remains
            self._y_limits = None   # trace set changed -> refit
            self._redraw()
            self.tracesChanged.emit()

    # -- loading -------------------------------------------------------------

    def add_dataset(self, path, label: str, dimensionality: int = 1,
                    pulse_program: str = "", nucleus: str = "") -> None:
        """Queue one dataset for loading and drawing."""
        path = Path(path)
        if any(t.path == path for t in self._traces):
            return   # already shown; dropping twice is a no-op, not a duplicate

        # Mixing dimensionalities in one figure is refused with an explicit
        # message. A 2D contour map and a 1D trace share no meaningful vertical
        # axis, so silently accepting the drop would produce a figure that
        # looks fine and means nothing.
        existing = [t for t in self._traces if t.visible]
        if existing:
            current_2d = any(t.is_2d for t in existing)
            incoming_2d = dimensionality == 2
            if current_2d != incoming_2d:
                self.dimensionalityRefused.emit(
                    f"Canvas is in {'2D' if current_2d else '1D'} mode. "
                    f"Clear it before adding a "
                    f"{'2D' if incoming_2d else '1D'} spectrum."
                )
                return
        self._pending_meta[Path(path)] = {
            "pulse_program": pulse_program, "nucleus": nucleus,
        }
        task = _LoadTask(self._reader, path, label, dimensionality)
        task.signals.loaded.connect(self._on_loaded)
        task.signals.loaded2d.connect(self._on_loaded_2d)
        task.signals.failed.connect(self._on_failed)
        # Hold a reference: QThreadPool does not keep the Python object alive,
        # and losing it mid-flight would silently discard the result.
        self._inflight.add(task)
        self._pool.start(task)

    def _on_loaded(self, path, ppm, intensity, label: str) -> None:
        path = Path(path)
        self._inflight = {t for t in self._inflight if t._path != path}
        if any(t.path == path for t in self._traces):
            return
        slot = self._slot_styles[len(self._traces) % len(self._slot_styles)]
        self._traces.append(
            Trace(
                path=path, label=label, ppm=ppm, intensity=intensity,
                color=slot["color"], line_width=slot["width"],
                line_style=slot["style"],
            )
        )
        meta = self._pending_meta.pop(path, {})
        self._traces[-1].pulse_program = meta.get("pulse_program", "")
        self._traces[-1].nucleus = meta.get("nucleus", "")
        self._y_limits = None   # new trace -> recompute the frame
        # Autoscale on every drop. Pinning the y range at the FIRST drop
        # (an earlier mistake) meant a spectrum loaded afterwards with a
        # much smaller intensity was drawn as a flat line at zero, because
        # the axis was still scaled to the first one. Limits are recomputed
        # whenever the set of traces changes, and only held fixed while the
        # user adjusts scale/offset -- which is what makes those
        # adjustments visible.
        self._y_limits = None
        # First spectrum dropped becomes the selection, so the wheel and
        # y-scale controls work immediately without an extra click.
        if self._selected_index is None:
            self._selected_index = len(self._traces) - 1
        self._redraw()
        self._emit_mode_if_changed()
        self.spectrumAdded.emit(label)
        self.tracesChanged.emit()

    def _on_loaded_2d(self, path, payload, label: str) -> None:
        path = Path(path)
        self._inflight = {t for t in self._inflight if t._path != path}
        if any(t.path == path for t in self._traces):
            return
        slot = self._slot_styles[len(self._traces) % len(self._slot_styles)]
        self._traces.append(
            Trace(
                path=path, label=label,
                ppm=np.asarray([]), intensity=np.asarray([]),
                color=slot["color"], line_width=slot["width"],
                line_style=slot["style"],
                is_2d=True,
                matrix=payload["matrix"],
                ppm_f1=payload["ppm_f1"],
                ppm_f2=payload["ppm_f2"],
            )
        )
        meta = self._pending_meta.pop(path, {})
        self._traces[-1].pulse_program = meta.get("pulse_program", "")
        self._traces[-1].nucleus = meta.get("nucleus", "")
        self._y_limits = None
        if self._selected_index is None:
            self._selected_index = len(self._traces) - 1
        self._redraw()
        self._emit_mode_if_changed()
        self.spectrumAdded.emit(label)
        self.tracesChanged.emit()

    def _on_failed(self, path, message: str) -> None:
        path = Path(path)
        self._inflight = {t for t in self._inflight if t._path != path}
        self.loadFailed.emit(str(path), message)

    # -- drawing -------------------------------------------------------------

    def _redraw(self) -> None:
        # 2D spectra get their own panel layout: contours cannot share an axis
        # with 1D traces meaningfully, and comparing 2D spectra means putting
        # them side by side.
        if any(t.visible and t.is_2d for t in self._traces):
            self._redraw_2d()
            return
        self._restore_single_axes()
        self._axes.clear()

        # Guard on VISIBLE traces, not merely on any traces existing: with
        # spectra loaded but all of them unchecked, the drawing code below has
        # nothing to take a max() over and previously raised ValueError,
        # crashing the redraw. Both "nothing loaded" and "all hidden" show the
        # empty state.
        visible_traces = [t for t in self._traces if t.visible]
        if not visible_traces:
            self._axes.set_xticks([])
            self._axes.set_yticks([])
            for spine in self._axes.spines.values():
                spine.set_visible(False)
            self._axes.text(
                0.5, 0.5,
                "Drag spectra here" if not self._traces
                else "All spectra hidden",
                ha="center", va="center", fontsize=13, color="#999999",
                transform=self._axes.transAxes,
            )
            self._canvas.draw_idle()
            return

        for spine in self._axes.spines.values():
            spine.set_visible(True)

        visible = visible_traces
        offset_step = 0.0
        if self._arrangement == self.ARRANGEMENT_STACKED and visible:
            # Measure the SCALED span: if a trace has been scaled up, the
            # stack spacing must grow with it or it overlaps its neighbour.
            tallest = max(
                (float(np.nanmax(t.intensity)) - float(np.nanmin(t.intensity)))
                * t.y_scale
                for t in visible
            )
            offset_step = tallest * 1.05 if tallest > 0 else 1.0

        selected = self.selected_trace()
        drawn = []
        for i, trace in enumerate(visible):
            # Per-trace vertical scale and offset, then the stacking offset.
            y = (trace.intensity * trace.y_scale) + trace.y_offset + (offset_step * i)
            width = trace.line_width
            if selected is not None and trace is selected and len(visible) > 1:
                width = width * 1.9   # make the selected trace obvious
            self._axes.plot(
                trace.ppm, y, color=trace.color, linewidth=width,
                linestyle=trace.line_style, label=trace.label,
            )
            drawn.append((trace, y))

        # ppm axes are conventionally DESCENDING (high ppm on the left).
        if self._ppm_range is not None:
            left, right = self._ppm_range
        else:
            left = max(float(np.nanmax(t.ppm)) for t in visible)
            right = min(float(np.nanmin(t.ppm)) for t in visible)
        self._axes.set_xlim(left, right)

        # Y limits come from the RAW data envelope (no per-trace scale or
        # offset applied), giving a STABLE frame that scaling is visible
        # against: scale a trace up and its peaks grow within the frame.
        #
        # An earlier version pinned the limits to whatever was first drawn and
        # never recomputed them. That made scaling visible but broke badly in
        # both directions -- scale down and the spectrum shrank to an invisible
        # line, scale up and it shot off the top with no way back. The frame is
        # now recomputed whenever the SET of traces changes (load, remove,
        # clear, arrangement) but deliberately NOT when scale or offset change.
        if self._y_limits is None:
            self._y_limits = self._frame_y_limits(visible, offset_step)
        if self._y_limits is not None:
            self._axes.set_ylim(*self._y_limits)

        self._axes.set_xlabel("ppm")
        if self._x_decimals is not None:
            from matplotlib.ticker import FormatStrFormatter

            self._axes.xaxis.set_major_formatter(
                FormatStrFormatter(f"%.{self._x_decimals}f")
            )
        if self._arrangement == self.ARRANGEMENT_STACKED:
            self._axes.set_yticks([])

        if self._show_grid:
            # Deliberately faint: a grid on a spectrum is a reading aid, not a
            # feature of the data, so it must not compete with the peaks.
            # X ONLY: a horizontal grid in a stacked plot cuts across every
            # spectrum and reads as part of the data. The ppm axis is the one
            # a reader actually measures against.
            if self._grid_spacing_ppm:
                from matplotlib.ticker import MultipleLocator

                self._axes.xaxis.set_major_locator(
                    MultipleLocator(self._grid_spacing_ppm)
                )
            self._axes.grid(
                True, axis="x", which="major",
                linewidth=0.4, alpha=0.25, linestyle="-",
            )
            self._axes.grid(False, axis="y")
        else:
            self._axes.grid(False)

        # Each spectrum's name sits at its own top-left, anchored in DATA
        # coordinates on x but pinned relative to its own trace on y, so the
        # label travels with its spectrum when stacked -- and is not dragged
        # around by y-scaling the way a legend entry would be.
        self._draw_trace_labels(drawn, left, right)

        self._figure.tight_layout()
        self._canvas.draw_idle()

    def _restore_single_axes(self) -> None:
        """Back to one axes, after a 2D multi-panel layout."""
        if len(self._figure.axes) != 1 or self._axes not in self._figure.axes:
            self._figure.clear()
            self._axes = self._figure.add_subplot(111)
            self._axes.set_facecolor("white")

    def _redraw_2d(self) -> None:
        """One panel per 2D spectrum, side by side, sharing the ppm axes.

        Side-by-side rather than overlaid: overlapping two contour maps is
        unreadable, which is why NMR software compares 2D spectra in adjacent
        panels. Shared axes mean zooming one moves all, so peaks stay aligned.
        """
        self._figure.clear()
        self._crosshair = None

        panels = [t for t in self._traces if t.visible and t.is_2d]
        one_d = [t for t in self._traces if t.visible and not t.is_2d]
        total = len(panels) + (1 if one_d else 0)
        if total == 0:
            self._axes = self._figure.add_subplot(111)
            self._axes.set_facecolor("white")
            self._canvas.draw_idle()
            return

        axes_list = self._figure.subplots(
            1, total, sharex=True, sharey=True, squeeze=False
        )[0]
        self._axes = axes_list[0]

        for ax, trace in zip(axes_list, panels):
            ax.set_facecolor("white")
            self._draw_contours(ax, trace)
            ax.set_title(
                trace.label, fontsize=8 * self._label_scale, color=trace.color
            )
            ax.set_xlabel("F2 (ppm)")
            f2 = self._f2_range or (
                float(np.nanmax(trace.ppm_f2)), float(np.nanmin(trace.ppm_f2))
            )
            f1 = self._f1_range or (
                float(np.nanmax(trace.ppm_f1)), float(np.nanmin(trace.ppm_f1))
            )
            ax.set_xlim(*f2)   # descending, NMR convention
            ax.set_ylim(*f1)
        axes_list[0].set_ylabel("F1 (ppm)")

        if one_d:
            ax = axes_list[-1]
            ax.set_facecolor("white")
            for trace in one_d:
                y = (trace.intensity * trace.y_scale) + trace.y_offset
                ax.plot(
                    trace.ppm, y, color=trace.color,
                    linewidth=trace.line_width,
                    linestyle=trace.line_style, label=trace.label,
                )
            ax.set_xlabel("ppm")

        self._figure.tight_layout()
        self._canvas.draw_idle()

    def _draw_contours(self, ax, trace) -> None:
        """Contour levels on a geometric ladder from a noise-based floor.

        NMR peaks span orders of magnitude, so evenly spaced levels either
        drown the plot in noise contours or show only the tallest peak. A
        geometric series starting just above the noise is what NMR displays
        use; y_scale doubles as the usual "contour level" control.
        """
        data = trace.matrix
        if data is None or data.size == 0:
            return
        finite = data[np.isfinite(data)]
        if finite.size == 0:
            return
        median = float(np.median(finite))
        mad = float(np.median(np.abs(finite - median))) * 1.4826
        base = (mad if mad > 0 else float(np.nanstd(finite))) * 4.0
        base = base / max(trace.y_scale, 1e-9)
        peak = float(np.nanmax(np.abs(finite)))
        if base <= 0 or peak <= base:
            base = peak / 20.0 if peak > 0 else 1.0
        n = max(2, int(trace.contour_levels))
        ratio = (peak / base) ** (1.0 / n) if base > 0 else 1.3
        levels = sorted({
            float(base * (ratio ** i)) for i in range(n)
            if np.isfinite(base * (ratio ** i))
        })
        if len(levels) < 2:
            return
        ax.contour(
            trace.ppm_f2, trace.ppm_f1, data,
            levels=levels, colors=trace.color,
            linewidths=trace.line_width, linestyles=trace.line_style,
        )

    def _frame_y_limits(self, visible, offset_step: float):
        """The stable vertical frame, from RAW intensities.

        Deliberately ignores each trace's y_scale and y_offset: those are what
        the user adjusts *within* the frame, so folding them in here would
        cancel out the very effect they are meant to have.
        """
        if not visible:
            return None
        lows, highs = [], []
        for i, trace in enumerate(visible):
            data = trace.intensity
            if data.size == 0:
                continue
            base = offset_step * i
            lows.append(float(np.nanmin(data)) + base)
            highs.append(float(np.nanmax(data)) + base)
        if not lows:
            return None
        low, high = min(lows), max(highs)
        if not np.isfinite(low) or not np.isfinite(high):
            return None
        if high == low:
            pad = 1.0 if high == 0 else abs(high) * 0.1
        else:
            pad = (high - low) * 0.08
        return (low - pad, high + pad)

    def _draw_trace_labels(self, drawn, left: float, right: float) -> None:
        """Name each spectrum at a FIXED position, top-left, one per line.

        Anchored in AXES-FRACTION coordinates, not data coordinates. An
        earlier version placed each label at its trace's data maximum, which
        meant scaling a spectrum dragged its label around the plot -- the
        reported bug. Fraction coordinates pin the labels to the corner of the
        axes regardless of scale, offset, or zoom.
        """
        if not drawn:
            return
        line_height = 0.045 * self._label_scale
        self._label_artists = []
        for i, (trace, _y) in enumerate(drawn):
            if trace.label_pos is None:
                base = (0.015, 0.985 - i * line_height)
                trace.label_base_pos = base
                dx, dy = trace.label_offset
                trace.label_pos = (
                    min(max(base[0] + dx, 0.0), 0.98),
                    min(max(base[1] + dy, 0.02), 1.0),
                )
            lx, ly = trace.label_pos
            lx += trace.label_dx
            ly += trace.label_dy
            lx = min(max(lx, -0.05), 1.05)
            ly = min(max(ly, -0.05), 1.05)
            artist = self._axes.text(
                lx, ly, trace.display_label(self._show_pulprog),
                color=trace.color, fontsize=8 * self._label_scale,
                ha="left", va="top",
                transform=self._axes.transAxes,   # <- fixed, not data coords
                clip_on=False,
            )
            self._label_artists.append((artist, trace))

    def autoscale_traces(self, target_fraction: float = 0.9) -> None:
        """Scale each spectrum individually so they are all actually visible.

        Without this, one spectrum three orders of magnitude stronger than
        another forces the weak one to a flat line -- the frame has to fit the
        strong one. Each trace is scaled so its own peak reaches the same
        fraction of the frame, which is what makes a comparison legible. The
        scales stay adjustable afterwards; this only sets a sane starting
        point.
        """
        visible = [t for t in self._traces if t.visible]
        if not visible:
            return
        frame = self._frame_y_limits(visible, 0.0)
        if frame is None:
            return
        span = frame[1] - frame[0]
        if span <= 0:
            return
        # In stacked mode each trace also gets an offset slot, so the share
        # of the frame available to any one of them is a fraction of the whole.
        divisor = len(visible) if self._arrangement == self.ARRANGEMENT_STACKED else 1
        target = (span * target_fraction) / divisor
        for trace in visible:
            peak = float(np.nanmax(np.abs(trace.intensity)))
            if peak > 0 and np.isfinite(peak):
                trace.y_scale = max(MIN_Y_SCALE, min(target / peak, MAX_Y_SCALE))
        self._redraw()
        # Then pin the frame around what is now actually drawn, so every
        # spectrum sits inside the canvas rather than running off the top.
        self.fit_to_drawn()
        self.tracesChanged.emit()

    def fit_to_drawn(self) -> None:
        """Frame exactly what is currently plotted, including scale and offset.

        Distinct from reset_y_limits(), which returns to the neutral raw-data
        frame. This one guarantees everything visible is inside the canvas --
        what "fit all spectra" means to a reader.
        """
        lows, highs = [], []
        for line in self._axes.lines:
            data = line.get_ydata()
            if data is None or len(data) == 0:
                continue
            arr = np.asarray(data, dtype=float)
            if not np.any(np.isfinite(arr)):
                continue
            lows.append(float(np.nanmin(arr)))
            highs.append(float(np.nanmax(arr)))
        if not lows:
            return
        low, high = min(lows), max(highs)
        if not np.isfinite(low) or not np.isfinite(high):
            return
        pad = (high - low) * 0.06 if high > low else (abs(high) * 0.1 or 1.0)
        self._y_limits = (low - pad, high + pad)
        self._axes.set_ylim(*self._y_limits)
        self._canvas.draw_idle()

    def move_to_bottom(self, index: int) -> None:
        """Sit one spectrum's baseline on the bottom of the frame."""
        if not (0 <= index < len(self._traces)):
            return
        trace = self._traces[index]
        frame = self._y_limits or self._frame_y_limits(
            [t for t in self._traces if t.visible], 0.0
        )
        if frame is None:
            return
        baseline = float(np.nanmin(trace.intensity)) * trace.y_scale
        trace.y_offset = frame[0] - baseline
        self._redraw()
        self.tracesChanged.emit()

    def move_all_to_bottom(self) -> None:
        """Every spectrum on a common baseline -- what 'overlay' means once
        offsets have been used. Equivalent to clearing every offset."""
        for i in range(len(self._traces)):
            self.move_to_bottom(i)

    # -- right-click export ---------------------------------------------

    # Format -> (file filter, is_vector). Vector formats keep text editable.
    EXPORT_FILTERS = [
        ("PNG image (*.png)", "png"),
        ("JPEG image (*.jpg)", "jpg"),
        ("TIFF image (*.tif)", "tif"),
        ("SVG vector (*.svg)", "svg"),
        ("PDF vector (*.pdf)", "pdf"),
        ("EPS vector (*.eps)", "eps"),
        ("PostScript (*.ps)", "ps"),
        ("PowerPoint slide (*.pptx)", "pptx"),
    ]
    # EMF is deliberately NOT offered. matplotlib cannot write it, and the
    # only route (converting via an external Inkscape install) silently
    # produced nothing when Inkscape was absent, which is worse than not
    # offering the format. SVG, PDF and EPS are vector and universally
    # readable; use one of those and convert downstream if EMF is required.

    def build_context_menu(self):
        """Menu construction only -- never calls exec().

        Split out so tests can trigger the actions without entering a real
        modal loop, which blocks forever in a test run.
        """
        menu = QMenu(self)
        menu.addAction("Save image\u2026", self.request_save_image)
        menu.addSeparator()
        menu.addAction("Auto scale Y", self.autoscale_traces)
        menu.addAction("Fit Y to data", self.fit_to_drawn)
        menu.addAction("All to bottom", self.move_all_to_bottom)
        return menu

    def _on_context_menu(self, pos) -> None:
        if not self._traces:
            return
        self.build_context_menu().exec(self.mapToGlobal(pos))

    def request_save_image(self) -> None:
        """Ask for a path and write the figure exactly as displayed."""
        filters = ";;".join(f for f, _ in self.EXPORT_FILTERS)
        path, chosen = QFileDialog.getSaveFileName(
            self, "Save spectrum image", "", filters
        )
        if not path:
            return
        # If the typed name has no extension, take it from the chosen filter
        # rather than silently writing a PNG named "figure".
        if "." not in path.rsplit("/", 1)[-1]:
            ext = next(
                (e for f, e in self.EXPORT_FILTERS if f == chosen), "png"
            )
            path = f"{path}.{ext}"
        try:
            self.save_image(path)
        except Exception as exc:  # noqa: BLE001 -- report, never crash
            self.loadFailed.emit(path, f"could not save: {exc}")
            return
        self.imageSaved.emit(path)

    def save_image(self, path, transparent: bool = False, dpi: int = 300) -> None:
        """Save exactly what is on screen.

        The tight-bounding-box save option is deliberately NOT used, and the
        layout is not altered: those silently re-crop and re-scale the figure,
        so the file would not match the canvas. The figure is written at its
        current size and limits, which is what "what you see is what you get"
        has to mean. (The banned keyword is enforced by test_architecture.py.)

        Vector formats (svg, pdf, eps, ps) keep text as editable text rather
        than outlines, so labels can be adjusted downstream.
        """
        import matplotlib as mpl

        # The crosshair is a UI overlay, not part of the figure. Leaving it in
        # bakes a stray cursor line into every exported image.
        self._clear_crosshair()

        path = str(path)
        suffix = path.rsplit(".", 1)[-1].lower() if "." in path else "png"
        facecolor = "none" if transparent else "white"

        # matplotlib converts text to outlines by default in vector output,
        # which makes labels uneditable in Illustrator/Inkscape -- the whole
        # point of exporting vector. 'none' for SVG keeps real <text>; type 42
        # (TrueType) for PDF/PS/EPS embeds editable text rather than Type 3.
        overrides = {
            "svg.fonttype": "none",
            "pdf.fonttype": 42,
            "ps.fonttype": 42,
        }
        if suffix == "pptx":
            self._save_pptx(path, dpi=dpi)
            return
        with mpl.rc_context(overrides):
            self._figure.savefig(
                path,
                format=suffix,
                dpi=dpi,
                facecolor=facecolor,
                edgecolor="none",
                transparent=transparent,
            )

    def _save_pptx(self, path: str, dpi: int = 300) -> None:
        """One slide sized to the figure, with the plot placed on it.

        The image is embedded as EMF-free PNG at full resolution: PowerPoint
        cannot render SVG reliably across versions, so a high-DPI raster is
        the format that actually survives being opened on someone else's
        machine.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:   # pragma: no cover - dependency present
            raise RuntimeError(
                "PowerPoint export needs the python-pptx package"
            ) from exc

        import tempfile as _tempfile

        width_in, height_in = self._figure.get_size_inches()
        with _tempfile.TemporaryDirectory() as tmp:
            png = f"{tmp}/slide.png"
            self._figure.savefig(
                png, format="png", dpi=dpi, facecolor="white", edgecolor="none"
            )
            prs = Presentation()
            prs.slide_width = Inches(float(width_in))
            prs.slide_height = Inches(float(height_in))
            # 6 = the blank layout in the default template.
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                png, 0, 0,
                width=Inches(float(width_in)),
                height=Inches(float(height_in)),
            )
            prs.save(path)

    def subtract(self, index_a: int, index_b: int, label: str | None = None) -> bool:
        """Add a difference spectrum (A - B) as a new trace.

        The two spectra are INTERPOLATED onto a common ppm axis before
        subtracting. Subtracting index-by-index is the classic silent error
        here: two spectra rarely have identical point counts or sweep widths,
        and index subtraction then quietly compares different chemical shifts
        and produces convincing nonsense.

        Their y_scale is applied first, so a difference taken after using
        "Same noise" or manual scaling reflects what is actually on screen.
        Returns False if the spectra do not overlap in ppm at all.
        """
        n = len(self._traces)
        if not (0 <= index_a < n and 0 <= index_b < n) or index_a == index_b:
            return False
        a, b = self._traces[index_a], self._traces[index_b]
        if a.is_2d or b.is_2d:
            return False   # 2D differences are not supported yet

        lo = max(float(np.nanmin(a.ppm)), float(np.nanmin(b.ppm)))
        hi = min(float(np.nanmax(a.ppm)), float(np.nanmax(b.ppm)))
        if not np.isfinite(lo) or not np.isfinite(hi) or hi <= lo:
            return False

        points = max(a.ppm.size, b.ppm.size)
        common = np.linspace(hi, lo, points)      # descending, as ppm axes are

        def resample(trace):
            # np.interp needs ascending x, so both are flipped and flipped back.
            order = np.argsort(trace.ppm)
            return np.interp(
                common[::-1],
                np.asarray(trace.ppm)[order],
                np.asarray(trace.intensity)[order] * trace.y_scale,
            )[::-1]

        difference = resample(a) - resample(b)
        # Mark it as derived so the legend and any later logic can tell a
        # difference apart from a measured spectrum.
        slot = self._slot_styles[len(self._traces) % len(self._slot_styles)]
        self._traces.append(
            Trace(
                path=Path(f"{a.path}::minus::{b.path}"),
                label=label or f"\u0394  {a.label}  \u2212  {b.label}",
                ppm=common,
                intensity=difference,
                color=slot["color"],
                line_width=slot["width"],
                line_style=slot["style"],
                is_difference=True,
            )
        )
        self._y_limits = None
        self._redraw()
        self.tracesChanged.emit()
        return True

    def set_label_offset(self, index: int, dx: float, dy: float) -> None:
        """Nudge one spectrum's name by an explicit amount, in axes fractions.

        Dragging is the quick way; this is the reliable, repeatable one --
        typed values survive a rebuild of the plot and can be dialled in
        precisely, which dragging on a dense spectrum cannot.
        """
        if not (0 <= index < len(self._traces)):
            return
        try:
            dx = float(dx)
            dy = float(dy)
        except (TypeError, ValueError):
            return
        if not (np.isfinite(dx) and np.isfinite(dy)):
            return
        self._traces[index].label_dx = dx
        self._traces[index].label_dy = dy
        self._redraw()

    def set_show_pulse_program(self, show: bool) -> None:
        self._show_pulprog = bool(show)
        self._redraw()

    def show_pulse_program(self) -> bool:
        return self._show_pulprog

    def _emit_mode_if_changed(self) -> None:
        current = self.mode()
        if current != self._last_mode:
            self._last_mode = current
            self.modeChanged.emit(current)

    def mode(self) -> str:
        """'2D' when any visible spectrum is 2D, otherwise '1D'.

        Derived from the data rather than being a separate switch the user has
        to remember to flip -- the app already knows each dataset's
        dimensionality.
        """
        return "2D" if any(t.visible and t.is_2d for t in self._traces) else "1D"

    def set_f1_range(self, high: float, low: float) -> None:
        """Indirect-dimension (F1) range, 2D only."""
        if high == low:
            return
        high, low = max(high, low), min(high, low)
        self._f1_range = (high, low)
        self._redraw()

    def set_f2_range(self, high: float, low: float) -> None:
        """Direct-dimension (F2) range, 2D only."""
        if high == low:
            return
        high, low = max(high, low), min(high, low)
        self._f2_range = (high, low)
        self._redraw()

    def f1_bounds(self):
        panels = [t for t in self._traces if t.visible and t.is_2d]
        if not panels:
            return None
        return (
            max(float(np.nanmax(t.ppm_f1)) for t in panels),
            min(float(np.nanmin(t.ppm_f1)) for t in panels),
        )

    def f2_bounds(self):
        panels = [t for t in self._traces if t.visible and t.is_2d]
        if not panels:
            return None
        return (
            max(float(np.nanmax(t.ppm_f2)) for t in panels),
            min(float(np.nanmin(t.ppm_f2)) for t in panels),
        )

    def display_label(self, trace) -> str:
        """Name shown on the plot: sample/expno plus the pulse programme.

        The pulse programme is what distinguishes two experiments on the same
        sample, which is exactly the case where a bare name is ambiguous.
        """
        if trace.pulse_program:
            return f"{trace.label}  ({trace.pulse_program})"
        return trace.label

    def set_label_offset(self, index: int, dx: float, dy: float) -> None:
        """Move one spectrum's name by an explicit amount, in axes fractions.

        Dragging is available too, but a typed offset is reproducible and
        works even when a name sits underneath a trace and is awkward to
        grab.
        """
        if not (0 <= index < len(self._traces)):
            return
        try:
            dx = float(dx)
            dy = float(dy)
        except (TypeError, ValueError):
            return
        if not (np.isfinite(dx) and np.isfinite(dy)):
            return
        trace = self._traces[index]
        base = trace.label_base_pos or trace.label_pos or (0.015, 0.985)
        trace.label_base_pos = base
        trace.label_offset = (dx, dy)
        trace.label_pos = (
            min(max(base[0] + dx, 0.0), 0.98),
            min(max(base[1] + dy, 0.02), 1.0),
        )
        self._redraw()
        self.tracesChanged.emit()

    def reset_label_positions(self) -> None:
        """Put dragged spectrum names back in the default stacked column."""
        for trace in self._traces:
            trace.label_pos = None
            trace.label_offset = (0.0, 0.0)
            trace.label_base_pos = None
            trace.label_dx = 0.0
            trace.label_dy = 0.0
        self._redraw()

    def set_label_scale(self, scale: float) -> None:
        """Size of the on-plot spectrum names, relative to default."""
        try:
            scale = float(scale)
        except (TypeError, ValueError):
            return
        if not np.isfinite(scale) or not (0.3 <= scale <= 4.0):
            return
        self._label_scale = scale
        self._redraw()

    def label_scale(self) -> float:
        return self._label_scale

    def set_x_decimals(self, decimals) -> None:
        """Decimal places on the ppm axis (0 -> "1", 2 -> "1.00"). None =
        let matplotlib choose."""
        if decimals is not None:
            try:
                decimals = int(decimals)
            except (TypeError, ValueError):
                return
            if decimals < 0 or decimals > 6:
                return
        self._x_decimals = decimals
        self._redraw()

    def x_decimals(self):
        return self._x_decimals

    def normalise_to_noise(self, target_snr_reference: int = 0) -> bool:
        """Scale every spectrum so their NOISE levels match.

        Equalising noise is what makes peak heights comparable between spectra
        acquired with different numbers of scans or receiver gain -- a peak
        twice as tall then genuinely means twice the signal-to-noise, which is
        the comparison a reader actually wants to make.

        Noise is estimated robustly as the median absolute deviation of the
        whole trace. MAD is used rather than the standard deviation because
        the peaks themselves are outliers: a plain sigma is dominated by the
        signal and would be a measure of peak height, not noise.

        Returns False when the noise cannot be estimated (e.g. a perfectly
        flat trace), rather than silently scaling by a meaningless factor.
        """
        visible = [t for t in self._traces if t.visible]
        if len(visible) < 1:
            return False
        if not (0 <= target_snr_reference < len(visible)):
            target_snr_reference = 0

        def noise_of(trace) -> float:
            data = np.asarray(trace.intensity, dtype=float)
            if data.size == 0:
                return 0.0
            median = float(np.nanmedian(data))
            mad = float(np.nanmedian(np.abs(data - median)))
            # 1.4826 converts MAD to a Gaussian-equivalent sigma.
            return mad * 1.4826

        noises = [noise_of(t) for t in visible]
        reference = noises[target_snr_reference]
        if reference <= 0 or not np.isfinite(reference):
            return False

        for trace, noise in zip(visible, noises):
            if noise <= 0 or not np.isfinite(noise):
                continue
            factor = reference / noise
            trace.y_scale = max(MIN_Y_SCALE, min(factor, MAX_Y_SCALE))
        self._redraw()
        self.fit_to_drawn()
        self.tracesChanged.emit()
        return True

    def set_grid_spacing_ppm(self, spacing) -> None:
        """Fixed grid spacing in ppm, or None to let matplotlib choose."""
        if spacing is not None:
            try:
                spacing = float(spacing)
            except (TypeError, ValueError):
                return
            if spacing <= 0 or not np.isfinite(spacing):
                return
        self._grid_spacing_ppm = spacing
        self._redraw()

    def grid_spacing_ppm(self):
        return self._grid_spacing_ppm

    def reset_y_limits(self) -> None:
        """Recompute the y range from what is currently drawn."""
        self._y_limits = None
        self._redraw()

    def apply_styles(self, styles: list[dict]) -> None:
        """Install per-slot appearance and re-style existing spectra.

        Slot i applies to the i-th loaded spectrum, so the mapping is the same
        before and after anything is dropped.
        """
        if not styles:
            return
        self._slot_styles = list(styles)
        for i, trace in enumerate(self._traces):
            slot = self._slot_styles[i % len(self._slot_styles)]
            trace.color = slot["color"]
            trace.line_style = slot["style"]
            trace.line_width = slot["width"]
        self._redraw()
        self.tracesChanged.emit()

    def slot_styles(self) -> list[dict]:
        return [dict(s) for s in self._slot_styles]

    def set_grid_visible(self, visible: bool) -> None:
        self._show_grid = bool(visible)
        self._redraw()

    def grid_visible(self) -> bool:
        return self._show_grid

    def set_trace_line_style(self, index: int, style: str) -> None:
        """Solid / dashed / dotted / dash-dot, per trace."""
        if not (0 <= index < len(self._traces)):
            return
        if style not in LINE_STYLES.values():
            return
        self._traces[index].line_style = style
        self._redraw()

    # -- cursor crosshair and drag-to-move -----------------------------------

    def set_crosshair_enabled(self, enabled: bool) -> None:
        self._crosshair_enabled = bool(enabled)
        if not enabled:
            self._clear_crosshair()

    def _clear_crosshair(self) -> None:
        if self._crosshair is not None:
            for artist in self._crosshair:
                try:
                    artist.remove()
                except (ValueError, NotImplementedError):
                    pass
            self._crosshair = None
            self._canvas.draw_idle()

    def _on_mouse_move(self, event) -> None:
        # Dragging the selected trace takes precedence over the crosshair.
        if self._label_drag is not None:
            self._continue_label_drag(event)
            return
        if self._drag_start is not None and event.inaxes is self._axes:
            self._continue_drag(event)
            return
        if event.inaxes is not self._axes:
            self._clear_crosshair()
            return
        if not getattr(self, "_crosshair_enabled", True):
            return
        self._clear_crosshair()
        vline = self._axes.axvline(
            event.xdata, color="#888888", linewidth=0.6, linestyle="--", alpha=0.8
        )
        hline = self._axes.axhline(
            event.ydata, color="#888888", linewidth=0.6, linestyle="--", alpha=0.8
        )
        # The ppm value at the cursor, drawn on the plot itself -- reading it
        # off the axis by eye is exactly what a crosshair is meant to avoid.
        label = self._axes.text(
            event.xdata, 1.005, f"{event.xdata:.4f} ppm",
            # x in DATA coords (follows the cursor), y in AXES coords (pinned
            # just above the frame) -- that is what get_xaxis_transform gives.
            transform=self._axes.get_xaxis_transform(),
            ha="center", va="bottom", fontsize=8, color="#333333",
            clip_on=False,
        )
        self._crosshair = (vline, hline, label)
        self.cursorMoved.emit(float(event.xdata), float(event.ydata))
        self._canvas.draw_idle()

    def _on_mouse_press(self, event) -> None:
        """Begin dragging the selected trace vertically.

        Only in stacked mode and only with a trace selected: dragging in
        overlay would be ambiguous about which spectrum is meant, and moving
        one without a clear selection would feel arbitrary.
        """
        if event.inaxes is not self._axes or event.button != 1:
            return
        # Dragging works in BOTH arrangements. Restricting it to stacked (an
        # earlier decision) meant the control silently did nothing in overlay,
        # which reads as broken rather than as a deliberate limitation. The
        # selection makes the target unambiguous either way.
        # A click on a spectrum NAME drags the name, not the spectrum -- the
        # names are what overlap the data and need moving out of the way.
        hit = self._label_at(event)
        if hit is not None:
            artist, hit_trace = hit
            self._label_drag = (hit_trace, event.x, event.y, hit_trace.label_pos)
            return

        trace = self.selected_trace()
        if trace is None or event.ydata is None:
            return
        self._drag_start = (float(event.ydata), trace.y_offset)

    def _label_at(self, event):
        """The spectrum name under the cursor, if any."""
        for artist, trace in getattr(self, "_label_artists", []):
            try:
                contains, _info = artist.contains(event)
            except (TypeError, AttributeError):
                continue
            if contains:
                return artist, trace
        return None

    def _continue_label_drag(self, event) -> None:
        trace, x0, y0, start_pos = self._label_drag
        if event.x is None or event.y is None:
            return
        width = max(self._canvas.width(), 1)
        height = max(self._canvas.height(), 1)
        dx = (event.x - x0) / width
        dy = (event.y - y0) / height
        nx = min(max(start_pos[0] + dx, 0.0), 0.98)
        ny = min(max(start_pos[1] + dy, 0.02), 1.0)
        trace.label_pos = (nx, ny)
        self._redraw()

    def _continue_drag(self, event) -> None:
        if event.ydata is None:
            return
        start_y, start_offset = self._drag_start
        index = self._selected_index
        if index is None:
            return
        self._traces[index].y_offset = start_offset + (float(event.ydata) - start_y)
        self._redraw()

    def _on_mouse_release(self, event) -> None:
        if self._label_drag is not None:
            self._label_drag = None
            self.tracesChanged.emit()
            return
        if self._drag_start is not None:
            self._drag_start = None
            self.tracesChanged.emit()

    def _on_mouse_leave(self, event) -> None:
        self._clear_crosshair()

    # -- drag and drop -------------------------------------------------------
    #
    # dragEnterEvent MUST call acceptProposedAction() or dropEvent never fires
    # -- the single most common Qt drag-and-drop mistake.

    def dragEnterEvent(self, event) -> None:
        if event.mimeData().hasFormat(MIME_DATASET):
            event.acceptProposedAction()

    def dragMoveEvent(self, event) -> None:
        if event.mimeData().hasFormat(MIME_DATASET):
            event.acceptProposedAction()

    def dropEvent(self, event) -> None:
        if self.handle_mime_data(event.mimeData()):
            event.acceptProposedAction()

    def handle_mime_data(self, mime) -> bool:
        """Decode a dataset payload and queue every entry.

        Separate from dropEvent so tests can exercise it with a plain
        QMimeData, without constructing QDropEvent objects whose signature
        varies between Qt versions.
        """
        if not mime.hasFormat(MIME_DATASET):
            return False
        try:
            payload = json.loads(bytes(mime.data(MIME_DATASET)).decode("utf-8"))
        except (ValueError, UnicodeDecodeError):
            return False
        if not payload:
            return False
        for item in payload:
            self.add_dataset(
                item["path"],
                item.get("label") or Path(item["path"]).name,
                int(item.get("dimensionality", 1)),
                pulse_program=item.get("pulse_program", ""),
                nucleus=item.get("nucleus", ""),
            )
        return True
