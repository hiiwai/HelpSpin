"""The main canvas: drop spectra straight onto it and see them.

Deliberately simple, replacing the earlier layout-first flow (define a figure
with N slots, then fill the slots). Here you just drag one or more datasets
from the browser onto the canvas and they are loaded and drawn. Arrangement
(overlay vs stacked) is a control you flip afterwards, not a decision you have
to make up front.

Loading happens on a worker thread: reading processed data off a network share
is slow enough to freeze the GUI if done inline. Each spectrum is drawn as it
arrives, so dropping several does not block on the slowest one.
"""

from __future__ import annotations

import json
import time
from dataclasses import dataclass
from pathlib import Path

import numpy as np
from matplotlib.backends.backend_qtagg import FigureCanvasQTAgg
from matplotlib.figure import Figure
from matplotlib.patches import Rectangle
from PySide6.QtCore import QObject, QRunnable, Qt, QThreadPool, Signal
from PySide6.QtWidgets import QFileDialog, QMenu, QVBoxLayout, QWidget

from ..domain.project import DEFAULT_PALETTE as PALETTE
from .dataset_model import MIME_DATASET

# Okabe-Ito is already the domain palette; reuse it so on-screen colours match
# whatever an eventual export produces.
# Yellow (#F0E442) is deliberately absent: it is nearly invisible as a thin
# line on white. Order puts high-contrast colours first so a two-spectrum
# comparison -- by far the common case -- is black against strong blue.
_FALLBACK_PALETTE = [
    "#000000", "#0072B2", "#D55E00", "#009E73",
    "#CC79A7", "#56B4E9", "#E69F00", "#8C564B",
]


# Real spectra in one figure can differ by many orders of magnitude (a 5 mM
# reference next to a 50 uM sample is ~1000x on its own), so the scale range
# has to span far more than a few decades. A 1000x ceiling silently clamped
# autoscale and made strong/weak pairs impossible to compare.
MIN_Y_SCALE = 1e-9
MAX_Y_SCALE = 1e9

LINE_STYLES = {
    "Solid": "-",
    "Dashed": "--",
    "Dotted": ":",
    "Dash-dot": "-.",
}


def _palette() -> list[str]:
    try:
        return list(PALETTE) or _FALLBACK_PALETTE
    except Exception:  # noqa: BLE001 - palette shape is not load-bearing here
        return _FALLBACK_PALETTE


class _PlotCanvas(FigureCanvasQTAgg):
    """Canvas subclass that guarantees mouse tracking.

    The wheel is handled through matplotlib's own ``scroll_event``, which the
    backend emits from FigureCanvasQT.wheelEvent. An earlier version tried to
    override that by assigning ``canvas.wheelEvent = fn`` on the INSTANCE --
    Qt dispatches virtuals through the class, so the assignment was silently
    ignored AND it shadowed the backend's own handler, killing scroll_event
    too. Hence: no wheel override here, just tracking so motion events arrive
    without a button held down.
    """

    def __init__(self, figure, owner):
        super().__init__(figure)
        self._owner = owner
        self.setMouseTracking(True)


@dataclass
class Trace:
    """One loaded 1D spectrum on the canvas.

    y_scale and y_offset are per-trace so a weak spectrum can be brought up to
    a strong one without touching the others -- adjusted with the scroll wheel
    over the plot, or typed exactly. line_width and color are per-trace too so
    a preference change can apply to all while still allowing overrides.
    """

    path: Path
    label: str
    ppm: np.ndarray
    intensity: np.ndarray
    color: str
    visible: bool = True
    y_scale: float = 1.0
    y_offset: float = 0.0
    # Horizontal shift in PPM, for lining up spectra that were referenced
    # slightly differently, or for skewing a stack into a cascade.
    #
    # This one is different in kind from y_offset and is treated with more
    # care because of it. Moving a trace vertically changes nothing a reader
    # would interpret; moving it horizontally moves it along the CHEMICAL
    # SHIFT axis, and chemical shift is data. A shifted trace is therefore
    # marked in the spectrum list and in the legend, so a figure never claims
    # a peak sits somewhere it does not. Peak readout always reports true,
    # unshifted ppm.
    x_offset: float = 0.0
    line_width: float = 0.8
    line_style: str = "-"      # matplotlib style: '-', '--', ':', '-.'
    # 2D spectra carry a matrix plus both ppm axes. ppm/intensity stay empty
    # for these; is_2d is what the drawing code branches on.
    is_2d: bool = False
    matrix: np.ndarray | None = None
    ppm_f1: np.ndarray | None = None
    ppm_f2: np.ndarray | None = None
    contour_levels: int = 12
    # Ratio between successive contour levels. NMR peaks span orders of
    # magnitude, so levels climb geometrically; this is the "multiplication
    # factor" a spectroscopist expects to set directly.
    contour_factor: float = 1.3
    # Multiple of the noise estimate for the LOWEST contour. Below ~3 the
    # plot fills with noise contours.
    contour_base_sigma: float = 4.0
    # Label position in AXES fractions, so dragging a name moves it and it
    # still does not drift when the spectrum is rescaled.
    label_pos: tuple[float, float] | None = None
    pulse_program: str = ""
    nucleus: str = ""
    # Number of scans and receiver gain. Bruker intensities scale linearly
    # with both, so two spectra acquired with different values are simply not
    # comparable as raw numbers -- NS 16/RG 101 against NS 512/RG 2050 is a
    # factor of ~640 before any chemistry is involved.
    ns: int = 0
    rg: float = 0.0
    label_offset: tuple[float, float] = (0.0, 0.0)
    label_base_pos: tuple[float, float] | None = None
    # Vestigial: a second, older label-offset mechanism whose setter was
    # shadowed by the label_offset one below and so never ran. Kept at 0.0
    # and still added when drawing, because the undo snapshot and saved
    # sessions from earlier versions carry the field; removing it would make
    # an old .helspin file fail to restore. Do not wire anything new to it.
    label_dx: float = 0.0     # always 0; see note above
    label_dy: float = 0.0
    is_difference: bool = False
    # How a derived trace was made. A difference has no file behind it, so
    # this recipe -- the two source paths, the operator, and the scales the
    # sources had AT THE TIME -- is what lets a session rebuild it. Recording
    # the scales matters: the arrays on screen were computed with the values
    # in force when Subtract was pressed, and re-deriving with whatever the
    # sources are scaled to later would silently produce a different result
    # from the one that was saved.
    source_a: Path | None = None
    source_b: Path | None = None
    operation: str = ""              # '-' or '+'
    source_scales: tuple[float, float] = (1.0, 1.0)

    def display_label(self, show_pulprog: bool = True) -> str:
        """Name as drawn on the plot.

        A non-zero x_offset is stated here, on the figure itself, and that is
        deliberate. A vertical offset is presentation; a horizontal one moves
        the trace along the chemical shift axis, so its peaks no longer read
        at their true values. A figure that has been aligned by hand should
        say so where the figure is looked at, not only in the application that
        made it.
        """
        name = self.label
        if show_pulprog and self.pulse_program:
            name = f"{name}  ({self.pulse_program})"
        if self.x_offset:
            name = f"{name}  [{self.x_offset:+.3f} ppm]"
        return name


class _LoadSignals(QObject):
    loaded = Signal(object, object, object, str, object)
    # path, ppm, intensity, label, (ns, rg) -- the acquisition
    # parameters have to cross with the data or they are lost: the
    # drag payload cannot know them, they only exist once the
    # spectrum has been read.
    loaded2d = Signal(object, object, str)         # path, payload dict, label
    failed = Signal(object, str)                   # path, message


def _payload_2d(spec) -> dict:
    """The contour payload, built in one place so the known-dimensionality
    path and the resolve-it-yourself path cannot diverge."""
    return {
        "is_2d": True,
        "matrix": np.asarray(spec.real, dtype=np.float64),
        "ppm_f1": np.asarray(spec.axis_f1.ppm_scale(), dtype=np.float64),
        "ppm_f2": np.asarray(spec.axis_f2.ppm_scale(), dtype=np.float64),
    }


class _LoadTask(QRunnable):
    """Reads one spectrum off the GUI thread.

    Pure I/O plus numpy: touches no Qt model state, so it is safe on a worker.
    The result is handed back by signal and applied on the GUI thread.
    """

    def __init__(self, reader, path: Path, label: str, dimensionality: int):
        super().__init__()
        # The canvas holds this task in _inflight until its result arrives.
        # Qt's default is to delete a QRunnable the moment run() returns, so
        # that reference would then point at freed memory and dropping it
        # would free it again -- a double free, i.e. a segfault with no
        # traceback. Python owns the task; Qt must not.
        self.setAutoDelete(False)
        self._reader = reader
        self._path = path
        self._label = label
        self._dim = dimensionality
        self.signals = _LoadSignals()

    def run(self) -> None:
        try:
            dim = self._dim
            if dim not in (1, 2):
                # 0 means the browser row was dropped before anything had
                # established whether it is 1D or 2D -- which is now the
                # normal case for a fresh row on a slow share, because rows
                # are draggable the moment they appear. Settling it here costs
                # a few stats on a worker thread; refusing the drop instead
                # (the old behaviour) cost the user the drag entirely.
                dim, spec = self._reader.read_auto(self._path)
                if dim == 2:
                    self.signals.loaded2d.emit(
                        self._path, _payload_2d(spec), self._label
                    )
                    return
                self.signals.loaded.emit(
                    self._path,
                    np.asarray(spec.axis.ppm_scale(), dtype=np.float64),
                    np.asarray(spec.real, dtype=np.float64),
                    self._label,
                    (getattr(spec, "ns", 0), getattr(spec, "rg", 0.0)),
                )
                return
            if dim == 2:
                spec = self._reader.read_2d(self._path)
                payload = {
                    "is_2d": True,
                    "matrix": np.asarray(spec.real, dtype=np.float64),
                    "ppm_f1": np.asarray(
                        spec.axis_f1.ppm_scale(), dtype=np.float64
                    ),
                    "ppm_f2": np.asarray(
                        spec.axis_f2.ppm_scale(), dtype=np.float64
                    ),
                }
                self.signals.loaded2d.emit(self._path, payload, self._label)
                return
            spec = self._reader.read_1d(self._path)
            ppm = np.asarray(spec.axis.ppm_scale(), dtype=np.float64)
            intensity = np.asarray(spec.real, dtype=np.float64)
            self.signals.loaded.emit(
                self._path, ppm, intensity, self._label,
                (getattr(spec, "ns", 0), getattr(spec, "rg", 0.0)),
            )
        except Exception as exc:  # noqa: BLE001 - one bad file must not kill the drop
            self.signals.failed.emit(self._path, str(exc))


def combine_arrays(ppm_a, int_a, scale_a, ppm_b, int_b, scale_b, op):
    """A op B on a common ppm axis. Returns (ppm, values) or None.

    Both spectra are INTERPOLATED onto a shared axis first. Subtracting
    index-by-index is the classic silent error here: two spectra rarely have
    identical point counts or sweep widths, so index subtraction quietly
    compares different chemical shifts and produces convincing nonsense.

    Module-level and array-only so the canvas and the session restore run the
    identical calculation -- a re-derived difference that did not match the
    one that was saved would be worse than not restoring it at all.
    """
    ppm_a = np.asarray(ppm_a, dtype=np.float64)
    ppm_b = np.asarray(ppm_b, dtype=np.float64)
    if ppm_a.size == 0 or ppm_b.size == 0:
        return None

    lo = max(float(np.nanmin(ppm_a)), float(np.nanmin(ppm_b)))
    hi = min(float(np.nanmax(ppm_a)), float(np.nanmax(ppm_b)))
    if not np.isfinite(lo) or not np.isfinite(hi) or hi <= lo:
        return None

    points = max(ppm_a.size, ppm_b.size)
    common = np.linspace(hi, lo, points)      # descending, as ppm axes are

    def resample(ppm, values, scale):
        # np.interp needs ascending x, so both are flipped and flipped back.
        order = np.argsort(ppm)
        return np.interp(
            common[::-1],
            np.asarray(ppm)[order],
            np.asarray(values, dtype=np.float64)[order] * scale,
        )[::-1]

    ya = resample(ppm_a, int_a, scale_a)
    yb = resample(ppm_b, int_b, scale_b)
    return common, (ya + yb if op == "+" else ya - yb)


class SpectrumCanvas(QWidget):
    """A matplotlib canvas that accepts dataset drops and plots them.

    Public surface kept small on purpose: drop things on it, call
    set_arrangement / set_ppm_range / clear. No figure-layout model, no slots.
    """

    spectrumAdded = Signal(str)      # label
    loadFailed = Signal(str, str)    # path, message
    tracesChanged = Signal()
    cursorMoved = Signal(float, float)    # ppm, intensity under the cursor
    imageSaved = Signal(str)              # path written
    modeChanged = Signal(str)             # '1D' or '2D'
    dimensionalityRefused = Signal(str)   # message
    viewChanged = Signal()                # ranges changed from the plot
    historyChanged = Signal()             # undo/redo availability changed

    ARRANGEMENT_OVERLAY = "overlay"
    ARRANGEMENT_STACKED = "stacked"

    def __init__(self, reader=None, pool: QThreadPool | None = None, parent=None):
        super().__init__(parent)
        if reader is None:
            from ..infrastructure.nmrglue_reader import NmrglueReader

            reader = NmrglueReader()
        self._reader = reader
        self._pool = pool or QThreadPool.globalInstance()
        self._traces: list[Trace] = []
        self._arrangement = self.ARRANGEMENT_OVERLAY
        self._ppm_range: tuple[float, float] | None = None
        self._inflight: set = set()
        self._selected_index: int | None = None
        self._default_line_width = 0.8
        self._palette_override: list[str] | None = None
        # Appearance per SLOT (1st spectrum loaded, 2nd, ...), set from
        # Preferences. Kept here so a newly dropped spectrum picks up its
        # slot's colour/style/width immediately.
        from .preferences_dialog import default_styles
        self._slot_styles = default_styles()
        # Y limits are held FIXED once established. Without this matplotlib
        # autoscales, so multiplying a trace's intensity just rescales the
        # axis and the spectrum looks completely unchanged -- the reported
        # 'Y scale does nothing' bug.
        self._y_limits: tuple[float, float] | None = None
        # The stacked lane height, remembered so that scaling one
        # spectrum does not re-lay the whole stack. Cleared wherever
        # _y_limits is, because the two describe one layout.
        self._stack_step: float | None = None
        self._zoom_mode = False
        self._y_zoom_mode = False
        # An EXPLICIT vertical window, set by the user with the Y zoom wheel
        # or a drag. Deliberately NOT stored in _y_limits, which is a cache of
        # the automatic frame and is cleared on almost every event (load,
        # remove, arrangement change, x zoom) to force a recompute. A user's
        # zoom must survive all of those, exactly as _ppm_range survives them
        # on the horizontal axis. Putting the two in one field would mean
        # every new spectrum silently threw the zoom away.
        self._y_range: tuple[float, float] | None = None
        # Display gain for STACKED mode only: a multiplier on every trace's
        # drawn amplitude, applied at draw time and never written into the
        # data. This is what the Y zoom wheel drives in a stack.
        #
        # Zooming the axis WINDOW is wrong in stacked mode. The lanes sit at
        # fixed positions, so narrowing the window about the cursor slides it
        # off the lower lanes and the bottom spectra vanish entirely -- the
        # reported fault. A stack is a fixed grid of lanes; zooming it means
        # making the traces taller in their lanes, not moving the lanes.
        self._stack_gain: float = 1.0
        self._zoom_band = None      # (x0, y0, Rectangle) while dragging
        self._opacity = 1.0
        # Digits after the point in the cursor readouts. Two is enough to
        # place a peak and keeps the numbers short; four filled the status
        # bar and the plot margin for no extra information.
        self._cursor_decimals = 2
        # Whether the spectrum names are drawn on the plot. On by default:
        # an overlay of unlabelled traces is unreadable. Off matters for a
        # figure whose caption names the spectra, and for a crowded 2D map
        # where the names sit over the data.
        self._labels_visible = True
        self._undo: list = []
        self._redo: list = []
        self._last_undo_key = None
        self._last_undo_at = 0.0
        self._show_grid = False
        # Grid spacing in ppm. None = let matplotlib choose.
        self._grid_spacing_ppm = None
        self._grid_spacing_y = None      # F1 in 2D, intensity in 1D
        # Decimal places on the ppm axis labels. None = matplotlib default.
        self._x_decimals = None
        # Relative size of the on-plot spectrum names. Larger values also
        # space them further apart, so a long list does not overlap the
        # traces it is labelling.
        self._label_scale = 1.0
        self._show_pulprog = True
        self._crosshair = None
        self._crosshair_enabled = True
        self._drag_start = None
        self._label_artists = []
        self._label_drag = None
        self._pending_meta = {}
        self._f1_range = None
        self._f2_range = None
        self._last_mode = "1D"

        self._figure = Figure(figsize=(6, 4), tight_layout=False, facecolor="white")
        self._canvas = _PlotCanvas(self._figure, self)
        self._axes = self._figure.add_subplot(111)
        self._axes.set_facecolor("white")

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(self._canvas)

        self.setAcceptDrops(True)
        self.setContextMenuPolicy(Qt.CustomContextMenu)
        self.customContextMenuRequested.connect(self._on_context_menu)
        # Scroll wheel over the plot scales the SELECTED trace vertically.
        self._canvas.setFocusPolicy(Qt.WheelFocus)
        # matplotlib's own scroll event rather than overriding Qt's
        # wheelEvent: assigning an instance attribute over a C++ virtual
        # is not reliably dispatched by the Shiboken binding.
        self._canvas.mpl_connect("scroll_event", self._on_scroll)
        # Crosshair + drag-to-move use matplotlib's own event system rather
        # than Qt's, so the coordinates arrive already in DATA space.
        self._canvas.mpl_connect("motion_notify_event", self._on_mouse_move)
        self._canvas.mpl_connect("button_press_event", self._on_mouse_press)
        self._canvas.mpl_connect("button_release_event", self._on_mouse_release)
        self._canvas.mpl_connect("axes_leave_event", self._on_mouse_leave)
        self._redraw()

    # -- public state --------------------------------------------------------

    @property
    def traces(self) -> list[Trace]:
        return list(self._traces)

    def arrangement(self) -> str:
        return self._arrangement

    def set_arrangement(self, arrangement: str) -> None:
        if arrangement not in (self.ARRANGEMENT_OVERLAY, self.ARRANGEMENT_STACKED):
            return
        self.push_undo()
        self._arrangement = arrangement
        self._y_limits = None   # stacking changes the frame entirely
        self._stack_step = None
        self._redraw()
        # Frame what the new arrangement actually drew. Switching layout is an
        # explicit "show me it this way", and the offsets that suit one layout
        # do not suit the other: after bottoming a stack, the offsets are
        # large, and the overlay frame -- which is built from raw data and
        # ignores offsets by design -- left the traces outside it. That is the
        # same disappearing-spectrum fault as the stacked one, in the other
        # direction.
        self.fit_to_drawn()

    def set_ppm_range(self, left: float, right: float) -> None:
        """ppm axes descend, so left must be the HIGHER value."""
        if left <= right:
            return
        self.push_undo("range")
        self._ppm_range = (left, right)
        # The vertical frame is derived from the data in view, so a new
        # horizontal window means a new vertical one. Without this, zooming
        # into a quiet region kept the scale set by peaks outside it and the
        # region of interest stayed a flat line.
        self._y_limits = None
        self._stack_step = None
        self._redraw()

    def full_range(self) -> None:
        """Show everything: the union of every trace's span, in every axis.

        The 2D ranges are cleared too. Only _ppm_range was reset before, so in
        2D mode -- where the view is driven by _f1_range and _f2_range -- the
        Full button cleared a value nothing was reading and appeared to do
        nothing at all. Union rather than intersection, so a spectrum with a
        wider sweep is shown in full rather than cropped to its neighbour.
        """
        self.push_undo()
        self._ppm_range = None
        self._f1_range = None
        self._f2_range = None
        self._y_limits = None   # the vertical frame follows the window
        self._stack_step = None
        self._redraw()
        self.viewChanged.emit()

    def _visible_1d(self) -> list:
        """Visible traces that actually carry 1D data.

        2D traces keep empty ppm/intensity arrays -- their data lives in
        matrix/ppm_f1/ppm_f2 -- so any 1D-only operation must filter them out.
        Passing an empty array to np.nanmax raises "zero-size array to
        reduction operation", which crashed the app the moment a 2D spectrum
        was dropped. Every 1D-only routine goes through this.
        """
        return [
            t for t in self._traces
            if t.visible and not t.is_2d and t.ppm.size and t.intensity.size
        ]

    def ppm_range(self):
        """The range currently displayed, or None when showing everything."""
        return self._ppm_range

    def f1_range(self):
        return self._f1_range

    def f2_range(self):
        return self._f2_range

    def ppm_bounds(self) -> tuple[float, float] | None:
        """Union of every visible trace's ppm span, or None if nothing loaded.

        Union rather than intersection: the range control should be able to
        show everything that exists, not only the overlap.
        """
        visible = self._visible_1d()
        if not visible:
            return None
        left = max(float(np.nanmax(t.ppm)) for t in visible)
        right = min(float(np.nanmin(t.ppm)) for t in visible)
        return (left, right)

    def clear(self) -> None:
        """Reset the canvas completely -- traces, selection, and ppm range.

        A true "start over", not just removing the lines: leaving a stale ppm
        range or selection behind after clearing is what makes the next drop
        appear on a nonsensical axis.
        """
        self.push_undo()   # clearing the canvas must be reversible
        self._traces.clear()
        self._ppm_range = None
        self._f1_range = None
        self._f2_range = None
        self._selected_index = None
        self._y_limits = None
        self._y_range = None    # an empty canvas has nothing to stay zoomed on
        self._stack_gain = 1.0
        self._redraw()
        self._emit_mode_if_changed()
        self.tracesChanged.emit()

    # -- selection and per-trace vertical scaling ----------------------------

    def selected_index(self) -> int | None:
        return self._selected_index

    def select_trace(self, index: int | None) -> None:
        """Select which trace the wheel / y-scale controls act on.

        None means 'nothing selected'; out-of-range indexes are ignored rather
        than raising, since a selection can outlive the trace it referred to.
        """
        if index is None:
            self._selected_index = None
        elif 0 <= index < len(self._traces):
            self._selected_index = index
        else:
            return
        # Deliberately does NOT emit tracesChanged: that rebuilds the spectra
        # list, which would wipe a multi-selection the moment a second
        # spectrum was picked -- exactly what stopped Subtract working.
        self._redraw()

    def selected_trace(self) -> Trace | None:
        if self._selected_index is None:
            return None
        if not (0 <= self._selected_index < len(self._traces)):
            return None
        return self._traces[self._selected_index]

    def set_y_scale(self, index: int, scale: float) -> None:
        """Set one trace's vertical scale. Non-positive values are refused --
        a zero or negative scale would flatten or invert the spectrum, which
        is never what a scale control is meant to do."""
        if not (0 <= index < len(self._traces)):
            return
        if scale <= 0 or not np.isfinite(scale):
            return
        scale = max(MIN_Y_SCALE, min(float(scale), MAX_Y_SCALE))
        self.push_undo(("scale", index))
        self._traces[index].y_scale = scale
        # NOT cleared for stacked mode. Re-fitting on every scale change made
        # turning one spectrum up shrink all the others, because the frame
        # grew to contain it. Holding the frame lets the scaled spectrum grow
        # into -- and past -- its lane and be clipped at the canvas edge,
        # leaving its neighbours exactly where they were. Fit Y re-fits.
        self._redraw()
        self.tracesChanged.emit()

    def nudge_y_scale(self, index: int, factor: float) -> None:
        """Multiply a trace's scale (wheel steps). Clamped to a sane range so
        a fast scroll cannot drive it to zero or to a value that overflows."""
        if not (0 <= index < len(self._traces)):
            return
        current = self._traces[index].y_scale
        new = current * float(factor)
        # Clamped to a range where the spectrum stays findable. The old
        # ceiling of 1e9 let a few scroll notches launch a trace far off
        # screen with no obvious way back.
        new = max(MIN_Y_SCALE, min(new, MAX_Y_SCALE))
        self.set_y_scale(index, new)

    def set_y_offset(self, index: int, offset: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        if not np.isfinite(offset):
            return
        self.push_undo(("offset", index))
        self._traces[index].y_offset = float(offset)
        # Offset moves the spectrum within the existing frame and NOTHING
        # else. It must not clear _y_limits: doing so recomputed the whole
        # frame from the lane positions, which rescaled the view so that every
        # spectrum appeared to resize when only one was moved. Offset is a
        # translation, not a zoom. The frame is fixed by the lane layout (see
        # _stacked_frame), which is deliberately built to hold offsets up to
        # the range the spin box allows, so a trace cannot be moved out of it.
        self._redraw()
        self.tracesChanged.emit()

    def nudge_y_offset(self, index: int, delta: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        self.set_y_offset(index, self._traces[index].y_offset + float(delta))

    def set_x_offset(self, index: int, offset: float) -> None:
        """Shift one spectrum along the ppm axis. Pure translation.

        The same rule as set_y_offset, for the same reason: this must NOT
        clear _ppm_range. Re-fitting the x limits to include the shifted trace
        would widen the frame every time one spectrum was nudged, so every
        OTHER spectrum would appear to slide and compress while the user was
        adjusting a single one. That regression was fixed once already on the
        y axis; it is exactly as wrong here.

        The consequence is deliberate: shift far enough and the trace runs off
        the edge of the view, rather than the view chasing it. The spin box is
        bounded to a fraction of the spectrum width so that cannot happen by
        accident.
        """
        if not (0 <= index < len(self._traces)):
            return
        if not np.isfinite(offset):
            return
        self.push_undo(("x_offset", index))
        self._traces[index].x_offset = float(offset)
        self._redraw()
        self.tracesChanged.emit()

    def nudge_x_offset(self, index: int, delta: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        self.set_x_offset(index, self._traces[index].x_offset + float(delta))

    def clear_x_offsets(self) -> None:
        """Return every spectrum to its true chemical shift.

        Offered as one action because that is how it is needed: having aligned
        several traces by eye, a user checking whether a difference is real
        wants them all back at once, not one at a time.
        """
        if not any(t.x_offset for t in self._traces):
            return
        self.push_undo(("x_offset_all", None))
        for trace in self._traces:
            trace.x_offset = 0.0
        self._redraw()
        self.tracesChanged.emit()

    def drawn_ppm(self, trace) -> np.ndarray:
        """The ppm axis AS DRAWN for this trace, shift included.

        A single place where the shift is applied, so the plot, the exports
        and the frame calculation cannot disagree about where a trace is.
        """
        if not trace.x_offset:
            return trace.ppm
        return np.asarray(trace.ppm, dtype=float) + float(trace.x_offset)

    def _on_scroll(self, event) -> None:
        """Wheel over the plot: zoom, or scale the selected trace.

        Three jobs on one wheel, resolved by the two zoom toggles:

        * either toggle on  -> zoom the view (both axes if both are on)
        * neither on        -> scale the SELECTED trace vertically

        The zoom checks come FIRST, and before the selection check, because
        zooming applies to the whole view and must not require a spectrum to
        be selected -- an earlier ordering would have made the new Y zoom
        silently dead until the user happened to click a trace.

        With no selection and no zoom mode the wheel does nothing, rather than
        scaling an arbitrary trace: guessing which spectrum was meant would be
        worse than requiring a click first.
        """
        step = getattr(event, "step", 0) or 0
        if step == 0:
            button = getattr(event, "button", None)
            step = 1 if button == "up" else (-1 if button == "down" else 0)
        if step == 0:
            return
        # ~10% per notch, direction following the scroll.
        factor = 1.1 if step > 0 else (1.0 / 1.1)
        if self._zoom_mode or self._y_zoom_mode:
            # Scrolling up magnifies, so the RANGE shrinks.
            self._zoom_about(event, 1.0 / factor)
            return
        index = self._selected_index
        if index is None:
            return
        self.nudge_y_scale(index, factor)

    # -- undo / redo ---------------------------------------------------------
    #
    # Snapshots record the STATE OF THE VIEW, not the data: the trace objects
    # themselves are held by reference, and only their mutable display fields
    # are copied. Undo therefore never re-reads a spectrum -- which rules out
    # the obvious alternative of reusing session_state()/restore_session(),
    # since that reloads every file from disk and would make undoing a scale
    # change cost a round trip per spectrum on a network share.
    #
    # Holding the removed Trace objects alive in the stack is also what makes
    # undoing a Remove instant and exact, data included.

    UNDO_DEPTH = 40

    # Space kept clear on the right for the crosshair's y-value label. It is
    # drawn OUTSIDE the axes, and tight_layout cannot know about an artist
    # added after it runs -- so without reserving this, the label was placed
    # correctly and then clipped off at the figure edge.
    RIGHT_MARGIN = 0.93

    def _snapshot(self) -> dict:
        return {
            "traces": [
                (t, t.y_scale, t.y_offset, t.x_offset, t.visible, t.color,
                 t.line_width, t.line_style, t.label_pos, t.label_offset)
                for t in self._traces
            ],
            "ppm": self._ppm_range,
            "f1": self._f1_range,
            "f2": self._f2_range,
            "arrangement": self._arrangement,
            "selected": self._selected_index,
            "y_limits": self._y_limits,
            "y_range": self._y_range,
            "stack_step": self._stack_step,
            "stack_gain": self._stack_gain,
        }

    def _apply_snapshot(self, snap: dict) -> None:
        self._traces = [entry[0] for entry in snap["traces"]]
        for (trace, scale, offset, x_offset, visible, color, width, style,
             label_pos, label_offset) in snap["traces"]:
            trace.y_scale, trace.y_offset, trace.visible = scale, offset, visible
            trace.x_offset = x_offset
            trace.color, trace.line_width, trace.line_style = color, width, style
            trace.label_pos, trace.label_offset = label_pos, label_offset
        self._ppm_range = snap["ppm"]
        self._f1_range = snap["f1"]
        self._f2_range = snap["f2"]
        self._arrangement = snap["arrangement"]
        self._selected_index = snap["selected"]
        self._y_limits = snap["y_limits"]
        # .get, not [], so a snapshot taken before this field existed (one
        # already on the undo stack when a session was restored) still applies
        # instead of raising KeyError halfway through an undo.
        self._y_range = snap.get("y_range")
        self._stack_step = snap["stack_step"]
        # .get with a default, for snapshots taken before this existed.
        self._stack_gain = snap.get("stack_gain", 1.0)
        self._redraw()
        self.tracesChanged.emit()
        self.viewChanged.emit()
        self.historyChanged.emit()

    # Bounds on the stacked display gain. Not decoration: without a floor a
    # few notches down flatten every trace to an invisible line, and without a
    # ceiling a few notches up put them so far past the frame that the canvas
    # shows nothing but vertical strokes -- both look like the spectra have
    # been lost, and both are reached in a couple of seconds of scrolling.

    COALESCE_SECONDS = 1.0

    def push_undo(self, key=None) -> None:
        """Record the current state before something changes it.

        ``key`` groups a burst of changes into ONE undo step. A spin box emits
        valueChanged on every keystroke, so typing "21.114" fired five
        separate changes and pushed five snapshots -- one undo then moved the
        scale from 21.114 to 21.11, which is indistinguishable from undo doing
        nothing at all. That was the reported fault. Repeated changes carrying
        the same key within a second are folded into the snapshot already
        taken, which is the state before the burst began, so one undo reverts
        the whole edit.

        Any new action clears the redo stack: once history has branched, the
        old forward path no longer describes anything reachable, and offering
        Redo for it would restore a state the user never returns to.
        """
        now = time.monotonic()
        if (key is not None and key == self._last_undo_key
                and now - self._last_undo_at < self.COALESCE_SECONDS):
            self._last_undo_at = now      # still typing: keep the burst open
            return
        self._last_undo_key, self._last_undo_at = key, now
        self._undo.append(self._snapshot())
        del self._undo[:-self.UNDO_DEPTH]
        if self._redo:
            self._redo.clear()
        self.historyChanged.emit()

    def can_undo(self) -> bool:
        return bool(self._undo)

    def can_redo(self) -> bool:
        return bool(self._redo)

    def undo(self) -> bool:
        if not self._undo:
            return False
        # End any open burst, so the next edit starts its own step rather than
        # being folded into the one just undone.
        self._last_undo_key = None
        self._redo.append(self._snapshot())
        self._apply_snapshot(self._undo.pop())
        return True

    def redo(self) -> bool:
        if not self._redo:
            return False
        self._undo.append(self._snapshot())
        self._apply_snapshot(self._redo.pop())
        return True

    # -- zooming -------------------------------------------------------------

    def set_zoom_mode(self, enabled: bool) -> None:
        """Wheel zooms about the cursor and dragging selects a box to zoom to.

        A mode rather than a modifier key: scaling the selected trace with the
        wheel and zooming with the wheel are both wanted often, and holding a
        key while scrolling is awkward. Off by default, so the wheel keeps the
        meaning it already had.
        """
        self._zoom_mode = bool(enabled)
        self._cancel_zoom_band()

    def zoom_mode(self) -> bool:
        return self._zoom_mode

    def set_y_zoom_mode(self, enabled: bool) -> None:
        """Wheel zooms the VERTICAL range, all spectra together.

        Independent of the horizontal toggle rather than exclusive with it:
        with both on the wheel zooms both axes at once, which is what the 2D
        view has always done and is the obvious reading of two separate
        switches. Turning this on does not disturb any zoom already applied.

        Distinct from the wheel's default job of scaling the SELECTED trace.
        That scales one spectrum against the others and changes their
        relationship; this magnifies the whole view and changes nothing
        between them.
        """
        self._y_zoom_mode = bool(enabled)

    def y_zoom_mode(self) -> bool:
        return self._y_zoom_mode

    def stack_gain(self) -> float:
        """Display magnification of the traces in stacked mode.

        1.0 means unzoomed. Lane positions are never affected by it.
        """
        return self._stack_gain

    def y_range(self) -> tuple[float, float] | None:
        """The explicit vertical window, or None when it is automatic."""
        return self._y_range

    def set_y_range(self, low: float, high: float) -> None:
        """Pin the vertical window explicitly.

        Rejects an inverted or degenerate pair rather than applying it: a
        zero-height axis renders as a blank plot with no obvious way back,
        which reads as the application having crashed.
        """
        try:
            low, high = float(low), float(high)
        except (TypeError, ValueError):
            return
        if not (np.isfinite(low) and np.isfinite(high)):
            return
        if high <= low:
            return
        self.push_undo("y_zoom")
        self._y_range = (low, high)
        self._redraw()
        self.viewChanged.emit()

    def reset_y_zoom(self) -> None:
        """Hand the vertical axis back to automatic framing."""
        if self._y_range is None and self._stack_gain == 1.0:
            return
        self.push_undo("y_zoom")
        self._y_range = None
        self._y_limits = None
        self._stack_gain = 1.0
        self._redraw()
        self.viewChanged.emit()

    def set_trace_opacity(self, opacity: float) -> None:
        """Alpha for every plotted line and contour.

        Superimposed 2D maps hide each other where they overlap, which is
        precisely where the interesting comparison is; partial opacity lets
        both be read at a crossing.
        """
        if not np.isfinite(opacity):
            return
        self._opacity = max(0.05, min(float(opacity), 1.0))
        self._redraw()

    def trace_opacity(self) -> float:
        return self._opacity

    def set_cursor_decimals(self, digits: int) -> None:
        self._cursor_decimals = max(0, min(int(digits), 6))

    def cursor_decimals(self) -> int:
        return self._cursor_decimals

    def _zoom_axis(self, low: float, high: float, centre: float,
                   factor: float) -> tuple:
        """Scale a range about `centre`. Order of the pair is preserved, so
        this works unchanged on a descending ppm axis."""
        return (
            centre + (low - centre) * factor,
            centre + (high - centre) * factor,
        )

    def _zoom_about(self, event, factor: float) -> None:
        """Zoom the visible range about the cursor.

        Keeping the point under the cursor fixed is the whole point: zooming
        about the centre of the plot walks the peak of interest off the edge
        and turns one gesture into three.
        """
        if event.xdata is None and event.ydata is None:
            return
        self.push_undo("zoom")
        if self.mode() == "2D":
            self._zoom_2d(event, factor)
        else:
            self._zoom_1d(event, factor)
        self._redraw()
        self.viewChanged.emit()

    def _zoom_2d(self, event, factor: float) -> None:
        """Both ppm axes. F1 is the vertical one here, so the Y toggle drives
        it -- but the horizontal toggle has always zoomed BOTH in 2D and that
        is left exactly as it was, or turning Y zoom off would silently remove
        behaviour people already rely on."""
        f2 = self._f2_range or self.f2_bounds()
        f1 = self._f1_range or self.f1_bounds()
        if self._zoom_mode and f2 is not None and event.xdata is not None:
            self._f2_range = self._zoom_axis(
                f2[0], f2[1], float(event.xdata), factor
            )
        if (self._zoom_mode or self._y_zoom_mode) and f1 is not None \
                and event.ydata is not None:
            self._f1_range = self._zoom_axis(
                f1[0], f1[1], float(event.ydata), factor
            )

    def _zoom_1d(self, event, factor: float) -> None:
        """The ppm window and/or the intensity window, per the two toggles."""
        if self._zoom_mode and event.xdata is not None:
            current = self._ppm_range or self.ppm_bounds()
            if current is not None:
                left, right = self._zoom_axis(
                    current[0], current[1], float(event.xdata), factor
                )
                # A descending axis that has collapsed or flipped would draw
                # an empty plot; refusing leaves the view as it was.
                if left > right:
                    self._ppm_range = (left, right)
                    # The horizontal window changing means the vertical
                    # AUTOMATIC frame must be recomputed -- but an explicit
                    # Y zoom is a stated intention and outranks that.
                    if self._y_range is None:
                        self._y_limits = None
                        self._stack_step = None
        if self._y_zoom_mode and self._arrangement == self.ARRANGEMENT_STACKED:
            # Stacked mode magnifies the traces instead of moving the window.
            # Nothing here touches _y_range, _y_limits or _stack_step, so the
            # lanes and the frame stay exactly where they are and no spectrum
            # can be zoomed out of sight.
            #
            # factor < 1 means "zoom in" (the caller inverts it), so the gain
            # is the reciprocal.
            if factor > 0:
                gain = self._stack_gain / factor
                # No ceiling and no floor. A weak signal beside a strong one
                # can need magnification of many orders of magnitude to read
                # at all, and a limit there just stops the wheel doing
                # anything with no explanation. The only rejection is a value
                # that is not a usable number: infinity or zero would draw a
                # blank plot with no way back, which is not a big zoom, it is
                # a broken one.
                if np.isfinite(gain) and gain > 0:
                    self._stack_gain = gain
            return
        if self._y_zoom_mode and event.ydata is not None:
            low, high = self._current_y_window()
            if low is not None:
                new_low, new_high = self._zoom_axis(
                    low, high, float(event.ydata), factor
                )
                if new_high > new_low and np.isfinite(new_low) \
                        and np.isfinite(new_high):
                    self._y_range = (new_low, new_high)

    def _current_y_window(self):
        """The vertical window in force, explicit or as drawn.

        Falls back to the axes' own limits so the first wheel notch zooms
        about what the user can actually see, rather than about a frame that
        has not been computed yet.
        """
        if self._y_range is not None:
            return self._y_range
        if self._y_limits is not None:
            return self._y_limits
        try:
            low, high = self._axes.get_ylim()
        except (AttributeError, ValueError):   # pragma: no cover
            return (None, None)
        if not (np.isfinite(low) and np.isfinite(high)) or high <= low:
            return (None, None)
        return (float(low), float(high))

    def _begin_zoom_band(self, event) -> bool:
        if event.xdata is None or event.ydata is None:
            return False
        rect = Rectangle(
            (float(event.xdata), float(event.ydata)), 0, 0,
            fill=False, edgecolor="#555555", linestyle="--", linewidth=0.8,
        )
        event.inaxes.add_patch(rect)
        self._zoom_band = (float(event.xdata), float(event.ydata), rect)
        return True

    def _update_zoom_band(self, event) -> None:
        if self._zoom_band is None or event.xdata is None or event.ydata is None:
            return
        x0, y0, rect = self._zoom_band
        rect.set_bounds(
            min(x0, float(event.xdata)), min(y0, float(event.ydata)),
            abs(float(event.xdata) - x0), abs(float(event.ydata) - y0),
        )
        self._canvas.draw_idle()

    def _cancel_zoom_band(self) -> None:
        if self._zoom_band is None:
            return
        _, _, rect = self._zoom_band
        try:
            rect.remove()
        except (ValueError, NotImplementedError):   # pragma: no cover
            pass
        self._zoom_band = None
        self._canvas.draw_idle()

    def _finish_zoom_band(self, event) -> None:
        """Apply the dragged box. A tiny drag is treated as a click.

        Without a minimum size, the smallest tremor while clicking would zoom
        to a few points of noise and there would be no obvious way back.
        """
        if self._zoom_band is None:
            return
        x0, y0, _ = self._zoom_band
        self._cancel_zoom_band()
        if event.xdata is None or event.ydata is None:
            return
        x1, y1 = float(event.xdata), float(event.ydata)
        if abs(x1 - x0) < 1e-9 or abs(y1 - y0) < 1e-12:
            return

        high_x, low_x = max(x0, x1), min(x0, x1)
        if self.mode() == "2D":
            self._f2_range = (high_x, low_x)
            self._f1_range = (max(y0, y1), min(y0, y1))
        else:
            self._ppm_range = (high_x, low_x)
            if self._y_zoom_mode:
                # With Y zoom on, a dragged box means the box: taking its
                # width but refitting its height would ignore half of a
                # gesture the user made deliberately.
                low_y, high_y = min(y0, y1), max(y0, y1)
                if high_y > low_y:
                    self._y_range = (low_y, high_y)
            else:
                self._y_range = None
                self._y_limits = None
                self._stack_step = None
        self._redraw()
        self.viewChanged.emit()

    # -- appearance preferences ---------------------------------------------

    def set_default_line_width(self, width: float) -> None:
        """Applies to every trace, including ones already drawn."""
        if width <= 0 or not np.isfinite(width):
            return
        self._default_line_width = float(width)
        for trace in self._traces:
            trace.line_width = float(width)
        self._redraw()

    def default_line_width(self) -> float:
        return self._default_line_width

    def set_trace_line_width(self, index: int, width: float) -> None:
        if not (0 <= index < len(self._traces)):
            return
        if width <= 0 or not np.isfinite(width):
            return
        self._traces[index].line_width = float(width)
        self._redraw()

    def set_trace_color(self, index: int, color: str) -> None:
        if not (0 <= index < len(self._traces)):
            return
        if not color:
            return
        self._traces[index].color = color
        self._redraw()
        self.tracesChanged.emit()

    def set_palette(self, colors: list[str]) -> None:
        """Replace the colour cycle and recolour existing traces in order."""
        if not colors:
            return
        self._palette_override = list(colors)
        for i, trace in enumerate(self._traces):
            trace.color = colors[i % len(colors)]
        self._redraw()
        self.tracesChanged.emit()

    def _active_palette(self) -> list[str]:
        return self._palette_override or _palette()

    def set_trace_visible(self, index: int, visible: bool) -> None:
        if not (0 <= index < len(self._traces)):
            return
        self._traces[index].visible = bool(visible)
        # Deliberately does NOT refit the frame. Hiding a spectrum is a
        # viewing action, not a change to the data being framed: refitting
        # made everything jump in scale the moment a box was unticked, and
        # re-ticking it did not restore the previous view. Only adding or
        # removing spectra, or switching arrangement, re-fits.
        self._redraw()
        self.tracesChanged.emit()

    def remove_trace(self, path: Path) -> None:
        # Snapshot BEFORE the removal, so the discarded Trace object stays
        # alive in the undo stack and coming back costs nothing -- no re-read
        # of the file, which on a share is the difference between instant and
        # a second per spectrum.
        snapshot = self._snapshot()
        before = len(self._traces)
        self._traces = [t for t in self._traces if t.path != Path(path)]
        if len(self._traces) != before:
            self._undo.append(snapshot)
            del self._undo[:-self.UNDO_DEPTH]
            self._redo.clear()
            self.historyChanged.emit()
            self._y_limits = None   # trace set changed -> refit
            self._redraw()
            self.tracesChanged.emit()

    # -- loading -------------------------------------------------------------

    def add_dataset(self, path, label: str, dimensionality: int = 1,
                    pulse_program: str = "", nucleus: str = "") -> None:
        """Queue one dataset for loading and drawing."""
        path = Path(path)
        if any(t.path == path for t in self._traces):
            return   # already shown; dropping twice is a no-op, not a duplicate

        # Mixing dimensionalities in one figure is refused with an explicit
        # message. A 2D contour map and a 1D trace share no meaningful vertical
        # axis, so silently accepting the drop would produce a figure that
        # looks fine and means nothing.
        # Only checked when the dimensionality is already known. A row
        # dropped before anything read it reports 0, and is treated exactly
        # like the several-at-once case below: dropping a 1D and a 2D in one
        # go is deliberately allowed and draws a panel each, so refusing this
        # one on arrival would contradict that.
        existing = [t for t in self._traces if t.visible]
        if existing and dimensionality in (1, 2):
            current_2d = any(t.is_2d for t in existing)
            incoming_2d = dimensionality == 2
            if current_2d != incoming_2d:
                self.dimensionalityRefused.emit(
                    f"Canvas is in {'2D' if current_2d else '1D'} mode. "
                    f"Clear it before adding a "
                    f"{'2D' if incoming_2d else '1D'} spectrum."
                )
                return
        self._pending_meta[Path(path)] = {
            "pulse_program": pulse_program, "nucleus": nucleus,
        }
        task = _LoadTask(self._reader, path, label, dimensionality)
        task.signals.loaded.connect(self._on_loaded)
        task.signals.loaded2d.connect(self._on_loaded_2d)
        task.signals.failed.connect(self._on_failed)
        # Hold a reference: QThreadPool does not keep the Python object alive,
        # and losing it mid-flight would silently discard the result.
        self._inflight.add(task)
        self._pool.start(task)

    def _on_loaded(self, path, ppm, intensity, label: str,
                   acquisition=None) -> None:
        path = Path(path)
        self._inflight = {t for t in self._inflight if t._path != path}
        if any(t.path == path for t in self._traces):
            return
        slot = self._slot_styles[len(self._traces) % len(self._slot_styles)]
        self._traces.append(
            Trace(
                path=path, label=label, ppm=ppm, intensity=intensity,
                color=slot["color"], line_width=slot["width"],
                line_style=slot["style"],
            )
        )
        meta = self._pending_meta.pop(path, {})
        self._traces[-1].pulse_program = meta.get("pulse_program", "")
        self._traces[-1].nucleus = meta.get("nucleus", "")
        if acquisition:
            self._traces[-1].ns = int(acquisition[0] or 0)
            self._traces[-1].rg = float(acquisition[1] or 0.0)
        self._y_limits = None   # new trace -> recompute the frame
        self._stack_step = None
        # Autoscale on every drop. Pinning the y range at the FIRST drop
        # (an earlier mistake) meant a spectrum loaded afterwards with a
        # much smaller intensity was drawn as a flat line at zero, because
        # the axis was still scaled to the first one. Limits are recomputed
        # whenever the set of traces changes, and only held fixed while the
        # user adjusts scale/offset -- which is what makes those
        # adjustments visible.
        self._y_limits = None
        # First spectrum dropped becomes the selection, so the wheel and
        # y-scale controls work immediately without an extra click.
        if self._selected_index is None:
            self._selected_index = len(self._traces) - 1
        self._redraw()
        self._emit_mode_if_changed()
        self.spectrumAdded.emit(label)
        self.tracesChanged.emit()

    def _on_loaded_2d(self, path, payload, label: str) -> None:
        path = Path(path)
        self._inflight = {t for t in self._inflight if t._path != path}
        if any(t.path == path for t in self._traces):
            return
        slot = self._slot_styles[len(self._traces) % len(self._slot_styles)]
        self._traces.append(
            Trace(
                path=path, label=label,
                ppm=np.asarray([]), intensity=np.asarray([]),
                color=slot["color"], line_width=slot["width"],
                line_style=slot["style"],
                is_2d=True,
                matrix=payload["matrix"],
                ppm_f1=payload["ppm_f1"],
                ppm_f2=payload["ppm_f2"],
            )
        )
        meta = self._pending_meta.pop(path, {})
        self._traces[-1].pulse_program = meta.get("pulse_program", "")
        self._traces[-1].nucleus = meta.get("nucleus", "")
        self._y_limits = None
        if self._selected_index is None:
            self._selected_index = len(self._traces) - 1
        self._redraw()
        self._emit_mode_if_changed()
        self.spectrumAdded.emit(label)
        self.tracesChanged.emit()

    def _on_failed(self, path, message: str) -> None:
        path = Path(path)
        self._inflight = {t for t in self._inflight if t._path != path}
        # The metadata stashed for this drop must go too. Leaving it meant a
        # failed load kept an entry keyed on that path for the rest of the
        # session, and the next successful load of the SAME dataset picked up
        # the stale PULPROG/nucleus from the attempt that failed.
        self._pending_meta.pop(path, None)
        self.loadFailed.emit(str(path), message)

    # -- drawing -------------------------------------------------------------

    def _redraw(self) -> None:
        # 2D spectra get their own panel layout: contours cannot share an axis
        # with 1D traces meaningfully, and comparing 2D spectra means putting
        # them side by side.
        if any(t.visible and t.is_2d for t in self._traces):
            self._redraw_2d()
            return
        self._restore_single_axes()
        self._axes.clear()

        # Guard on VISIBLE traces, not merely on any traces existing: with
        # spectra loaded but all of them unchecked, the drawing code below has
        # nothing to take a max() over and previously raised ValueError,
        # crashing the redraw. Both "nothing loaded" and "all hidden" show the
        # empty state.
        visible_traces = [t for t in self._traces if t.visible]
        if not visible_traces:
            self._axes.set_xticks([])
            self._axes.set_yticks([])
            for spine in self._axes.spines.values():
                spine.set_visible(False)
            self._axes.text(
                0.5, 0.5,
                "Drag spectra here" if not self._traces
                else "All spectra hidden",
                ha="center", va="center", fontsize=13, color="#999999",
                transform=self._axes.transAxes,
            )
            self._canvas.draw_idle()
            return

        for spine in self._axes.spines.values():
            spine.set_visible(True)

        visible = visible_traces
        offset_step = 0.0
        if self._arrangement == self.ARRANGEMENT_STACKED and visible:
            # Measure the SCALED span: if a trace has been scaled up, the
            # stack spacing must grow with it or it overlaps its neighbour.
            # Measured over the VISIBLE window: a peak outside the view has
            # no business deciding how tall a lane is on screen.
            if self._stack_step is None:
                spans = []
                for t in visible:
                    seen = self._window_values(t)
                    if seen.size:
                        spans.append(
                            (float(np.nanmax(seen)) - float(np.nanmin(seen)))
                            * t.y_scale
                        )
                tallest = max(spans) if spans else 0.0
                self._stack_step = tallest * 1.05 if tallest > 0 else 1.0
            # The lane grid is REMEMBERED, not recomputed on every redraw.
            # Deriving it from the tallest scaled span each time meant turning
            # one spectrum up moved every other spectrum: scaling the middle
            # of three by 20 grew the frame 13-fold and squashed its
            # neighbours from 29% of the canvas to 2%. A stack is a fixed set
            # of lanes; scaling a spectrum makes it taller within its lane and,
            # past a point, clipped by the canvas -- which is what a stacked
            # display is expected to do. Fit Y re-establishes the grid.
            offset_step = self._stack_step

        selected = self.selected_trace()
        drawn = []
        for i, trace in enumerate(visible):
            # Per-trace vertical scale and offset, then the stacking offset.
            y = ((trace.intensity * trace.y_scale * self._stack_gain)
                 + trace.y_offset + (offset_step * i))
            width = trace.line_width
            if selected is not None and trace is selected and len(visible) > 1:
                width = width * 1.9   # make the selected trace obvious
            self._axes.plot(
                self.drawn_ppm(trace), y, color=trace.color, linewidth=width,
                linestyle=trace.line_style, label=trace.label,
                alpha=self._opacity,
            )
            drawn.append((trace, y))

        # ppm axes are conventionally DESCENDING (high ppm on the left).
        if self._ppm_range is not None:
            left, right = self._ppm_range
        else:
            left = max(float(np.nanmax(t.ppm)) for t in visible)
            right = min(float(np.nanmin(t.ppm)) for t in visible)
        self._axes.set_xlim(left, right)

        # Y limits come from the RAW data envelope (no per-trace scale or
        # offset applied), giving a STABLE frame that scaling is visible
        # against: scale a trace up and its peaks grow within the frame.
        #
        # An earlier version pinned the limits to whatever was first drawn and
        # never recomputed them. That made scaling visible but broke badly in
        # both directions -- scale down and the spectrum shrank to an invisible
        # line, scale up and it shot off the top with no way back. The frame is
        # now recomputed whenever the SET of traces changes (load, remove,
        # clear, arrangement) but deliberately NOT when scale or offset change.
        # An explicit Y zoom outranks every automatic frame. Checked before
        # the cache so that the recompute below is skipped entirely: running
        # it and then overriding would waste the work and, in stacked mode,
        # would also re-lay the lanes underneath a user who had zoomed in.
        if self._y_range is not None:
            self._axes.set_ylim(*self._y_range)
        elif self._y_limits is None:
            if self._arrangement == self.ARRANGEMENT_STACKED:
                # Stacked framing has to use the DRAWN positions. The overlay
                # rule -- frame the raw envelope, ignore scale and offset --
                # cannot work here, because the lane spacing is itself a
                # SCALED quantity. Mixing the two put the frame on one scale
                # and the traces on another: with y_scale 21 and y_offset
                # -2.04e12 the frame came out as -1.9e11..2.5e12 while the
                # lower trace was drawn at -2.04e12..7.1e10, entirely below
                # it. That is the "one spectrum is gone in stacked" report.
                self._y_limits = self._stacked_frame(visible, offset_step)
            else:
                self._y_limits = self._frame_y_limits(visible, offset_step)
            if self._y_limits is not None:
                self._axes.set_ylim(*self._y_limits)
        elif self._y_limits is not None:
            self._axes.set_ylim(*self._y_limits)

        self._axes.set_xlabel("ppm")
        if self._x_decimals is not None:
            from matplotlib.ticker import FormatStrFormatter

            self._axes.xaxis.set_major_formatter(
                FormatStrFormatter(f"%.{self._x_decimals}f")
            )
        if self._arrangement == self.ARRANGEMENT_STACKED:
            self._axes.set_yticks([])

        if self._show_grid:
            self._apply_grid(
                self._axes,
                # A horizontal grid is drawn in OVERLAY only. Stacked mode
                # gives every spectrum its own offset and scale, so a y value
                # means something different for each one; ruling lines across
                # them would imply a shared scale that does not exist. Say so
                # rather than drawing something misleading.
                horizontal=self._arrangement != self.ARRANGEMENT_STACKED,
            )
        else:
            self._axes.grid(False)

        # Each spectrum's name sits at its own top-left, anchored in DATA
        # coordinates on x but pinned relative to its own trace on y, so the
        # label travels with its spectrum when stacked -- and is not dragged
        # around by y-scaling the way a legend entry would be.
        self._draw_trace_labels(drawn, left, right)

        self._figure.tight_layout(rect=(0, 0, self.RIGHT_MARGIN, 1))
        self._canvas.draw_idle()

    def _restore_single_axes(self) -> None:
        """Back to one axes, after a 2D multi-panel layout."""
        if len(self._figure.axes) != 1 or self._axes not in self._figure.axes:
            self._figure.clear()
            self._axes = self._figure.add_subplot(111)
            self._axes.set_facecolor("white")

    def _redraw_2d(self) -> None:
        """2D layout, following the SAME arrangement setting as 1D.

        Overlay superimposes the contour maps on one set of axes, each in its
        own colour; Stacked puts them in adjacent panels. Until 0.4.9 both
        settings drew panels, so "Overlay" and "Stacked" were the same thing
        and there was no way to superimpose at all -- which is the comparison
        that matters most in 2D. Chemical-shift perturbation work (apo against
        ligand-bound, the case in the report) is read by looking at how far
        each peak MOVED, and that only shows when the two maps share axes.
        Adjacent panels are still the right answer for maps too crowded to
        superimpose, which is why both exist.
        """
        if self._arrangement == self.ARRANGEMENT_OVERLAY:
            self._redraw_2d_overlay()
            return
        self._redraw_2d_panels()

    def _redraw_2d_overlay(self) -> None:
        """Every 2D map on one set of axes, in its own colour."""
        self._figure.clear()
        self._crosshair = None
        maps = [t for t in self._traces if t.visible and t.is_2d]
        one_d = [t for t in self._traces if t.visible and not t.is_2d]

        if one_d:
            # A 1D trace cannot share a vertical axis with a contour map: one
            # is intensity, the other is ppm. It keeps its own panel.
            axes_list = self._figure.subplots(1, 2, squeeze=False)[0]
            self._axes, one_d_axes = axes_list[0], axes_list[1]
        else:
            self._axes = self._figure.add_subplot(111)
            one_d_axes = None

        self._axes.set_facecolor("white")
        for trace in maps:
            self._draw_contours(self._axes, trace)

        if maps:
            f2 = self._f2_range or (
                max(float(np.nanmax(t.ppm_f2)) for t in maps),
                min(float(np.nanmin(t.ppm_f2)) for t in maps),
            )
            f1 = self._f1_range or (
                max(float(np.nanmax(t.ppm_f1)) for t in maps),
                min(float(np.nanmin(t.ppm_f1)) for t in maps),
            )
            self._axes.set_xlim(*f2)   # descending, NMR convention
            self._axes.set_ylim(*f1)
        self._axes.set_xlabel("F2 (ppm)")
        self._axes.set_ylabel("F1 (ppm)")
        if self._show_grid:
            # Both axes here: in 2D they are both chemical shifts, and reading
            # a peak position off a contour map is exactly what a grid helps
            # with.
            self._apply_grid(self._axes)

        # Superimposed maps have no titles to tell them apart, so each is
        # named in the corner in its own colour -- the same device the 1D
        # overlay uses, and for the same reason.
        for row, trace in enumerate(maps if self._labels_visible else []):
            self._axes.text(
                0.01, 0.99 - row * 0.045, trace.label,
                transform=self._axes.transAxes,
                ha="left", va="top", color=trace.color,
                fontsize=8 * self._label_scale,
            )

        if one_d_axes is not None:
            one_d_axes.set_facecolor("white")
            for trace in one_d:
                y = (trace.intensity * trace.y_scale) + trace.y_offset
                one_d_axes.plot(
                    self.drawn_ppm(trace), y, color=trace.color,
                    linewidth=trace.line_width, linestyle=trace.line_style,
                    label=trace.label,
                )
            one_d_axes.set_xlabel("ppm")

        self._figure.tight_layout(rect=(0, 0, self.RIGHT_MARGIN, 1))
        self._canvas.draw_idle()

    def _redraw_2d_panels(self) -> None:
        """One panel per 2D spectrum, side by side, sharing the ppm axes.

        The right answer when the maps are too crowded to superimpose. Shared
        axes mean zooming one moves all, so peaks stay aligned.
        """
        self._figure.clear()
        self._crosshair = None

        panels = [t for t in self._traces if t.visible and t.is_2d]
        one_d = [t for t in self._traces if t.visible and not t.is_2d]
        total = len(panels) + (1 if one_d else 0)
        if total == 0:
            self._axes = self._figure.add_subplot(111)
            self._axes.set_facecolor("white")
            self._canvas.draw_idle()
            return

        axes_list = self._figure.subplots(
            1, total, sharex=True, sharey=True, squeeze=False
        )[0]
        self._axes = axes_list[0]

        for ax, trace in zip(axes_list, panels):
            ax.set_facecolor("white")
            self._draw_contours(ax, trace)
            ax.set_title(
                trace.label, fontsize=8 * self._label_scale, color=trace.color
            )
            ax.set_xlabel("F2 (ppm)")
            if self._show_grid:
                self._apply_grid(ax)
            f2 = self._f2_range or (
                float(np.nanmax(trace.ppm_f2)), float(np.nanmin(trace.ppm_f2))
            )
            f1 = self._f1_range or (
                float(np.nanmax(trace.ppm_f1)), float(np.nanmin(trace.ppm_f1))
            )
            ax.set_xlim(*f2)   # descending, NMR convention
            ax.set_ylim(*f1)
        axes_list[0].set_ylabel("F1 (ppm)")

        if one_d:
            ax = axes_list[-1]
            ax.set_facecolor("white")
            for trace in one_d:
                y = (trace.intensity * trace.y_scale) + trace.y_offset
                ax.plot(
                    self.drawn_ppm(trace), y, color=trace.color,
                    linewidth=trace.line_width,
                    linestyle=trace.line_style, label=trace.label,
                )
            ax.set_xlabel("ppm")

        self._figure.tight_layout(rect=(0, 0, self.RIGHT_MARGIN, 1))
        self._canvas.draw_idle()

    def _draw_contours(self, ax, trace) -> None:
        """Contour levels on a geometric ladder from a noise-based floor.

        NMR peaks span orders of magnitude, so evenly spaced levels either
        drown the plot in noise contours or show only the tallest peak. A
        geometric series starting just above the noise is what NMR displays
        use; y_scale doubles as the usual "contour level" control.
        """
        data = trace.matrix
        if data is None or data.size == 0:
            return
        finite = data[np.isfinite(data)]
        if finite.size == 0:
            return
        median = float(np.median(finite))
        mad = float(np.median(np.abs(finite - median))) * 1.4826
        sigma = mad if mad > 0 else float(np.nanstd(finite))
        base = sigma * float(trace.contour_base_sigma or 4.0)
        base = base / max(trace.y_scale, 1e-9)
        peak = float(np.nanmax(np.abs(finite)))
        if base <= 0 or peak <= base:
            base = peak / 20.0 if peak > 0 else 1.0
        n = max(2, int(trace.contour_levels))
        ratio = float(trace.contour_factor)
        if not np.isfinite(ratio) or ratio <= 1.0:
            # A factor of 1 or less would repeat the same level forever.
            ratio = 1.3
        levels = sorted({
            float(base * (ratio ** i)) for i in range(n)
            if np.isfinite(base * (ratio ** i))
        })
        if len(levels) < 2:
            return
        ax.contour(
            trace.ppm_f2, trace.ppm_f1, data,
            levels=levels, colors=trace.color,
            linewidths=trace.line_width, linestyles=trace.line_style,
            alpha=self._opacity,
        )

    def _stacked_frame(self, visible, offset_step: float):
        """The frame for stacked mode: the union of the drawn lanes.

        Every term that positions a trace on screen is included -- y_scale,
        y_offset and the stack step -- so a lane cannot fall outside the
        frame. That is the whole contract of stacked mode: N spectra, N lanes,
        all of them on the canvas.

        With the step set to the tallest scaled span, the lanes come out
        evenly spaced, so "to bottom" on the second of two spectra puts its
        baseline halfway up, on the third of three a third of the way up, and
        so on -- which is what a stack is.
        """
        lows, highs = [], []
        for i, trace in enumerate(visible):
            data = self._window_values(trace)
            if data.size == 0:
                continue
            base = trace.y_offset + offset_step * i
            # The Y-zoom gain is deliberately NOT included.
            #
            # It used to be, on the reasoning that no lane should fall outside
            # the frame. But this frame is recomputed on load, remove and
            # window change, so including the gain meant the frame grew to
            # swallow the magnification the moment any of those happened:
            # zoom in, drop in one more spectrum, and the peaks silently
            # shrank back towards their original size. The magnification
            # appeared to undo itself for no reason the user could see.
            #
            # Excluding it makes the frame a property of the LANES alone --
            # baselines and the unmagnified envelope -- so it does not move
            # when the gain changes, and the gain is a stable magnification
            # rather than one the next load renegotiates. Magnified peaks
            # overflow their lane and can run off the canvas entirely, which
            # is intended: that is what turning the intensity up in a stacked
            # plot does, and it is the caller's business how far to go.
            scale = trace.y_scale
            lows.append(float(np.nanmin(data)) * scale + base)
            highs.append(float(np.nanmax(data)) * scale + base)
        if not lows:
            return None
        low, high = min(lows), max(highs)
        if not np.isfinite(low) or not np.isfinite(high):
            return None
        if high == low:
            pad = 1.0 if high == 0 else abs(high) * 0.1
        else:
            pad = (high - low) * 0.04
        return (low - pad, high + pad)

    def _frame_y_limits(self, visible, offset_step: float):
        """The stable vertical frame, from RAW intensities.

        Deliberately ignores each trace's y_scale and y_offset: those are what
        the user adjusts *within* the frame, so folding them in here would
        cancel out the very effect they are meant to have.
        """
        if not visible:
            return None
        lows, highs = [], []
        for i, trace in enumerate(visible):
            data = self._window_values(trace)
            if data.size == 0:
                continue
            base = offset_step * i
            lows.append(float(np.nanmin(data)) + base)
            highs.append(float(np.nanmax(data)) + base)
        if not lows:
            return None
        low, high = min(lows), max(highs)
        if not np.isfinite(low) or not np.isfinite(high):
            return None
        if high == low:
            pad = 1.0 if high == 0 else abs(high) * 0.1
        else:
            pad = (high - low) * 0.08
        return (low - pad, high + pad)

    def set_labels_visible(self, visible: bool) -> None:
        self._labels_visible = bool(visible)
        self._redraw()

    def labels_visible(self) -> bool:
        return self._labels_visible

    def _draw_trace_labels(self, drawn, left: float, right: float) -> None:
        """Name each spectrum at a FIXED position, top-left, one per line.

        Anchored in AXES-FRACTION coordinates, not data coordinates. An
        earlier version placed each label at its trace's data maximum, which
        meant scaling a spectrum dragged its label around the plot -- the
        reported bug. Fraction coordinates pin the labels to the corner of the
        axes regardless of scale, offset, or zoom.
        """
        if not drawn or not self._labels_visible:
            return
        line_height = 0.045 * self._label_scale
        self._label_artists = []
        for i, (trace, _y) in enumerate(drawn):
            if trace.label_pos is None:
                base = (0.015, 0.985 - i * line_height)
                trace.label_base_pos = base
                dx, dy = trace.label_offset
                trace.label_pos = (
                    min(max(base[0] + dx, 0.0), 0.98),
                    min(max(base[1] + dy, 0.02), 1.0),
                )
            lx, ly = trace.label_pos
            lx += trace.label_dx
            ly += trace.label_dy
            lx = min(max(lx, -0.05), 1.05)
            ly = min(max(ly, -0.05), 1.05)
            artist = self._axes.text(
                lx, ly, trace.display_label(self._show_pulprog),
                color=trace.color, fontsize=8 * self._label_scale,
                ha="left", va="top",
                transform=self._axes.transAxes,   # <- fixed, not data coords
                clip_on=False,
            )
            self._label_artists.append((artist, trace))

    def autoscale_traces(self, target_fraction: float = 0.9) -> None:
        """Scale each spectrum individually so they are all actually visible.

        Without this, one spectrum three orders of magnitude stronger than
        another forces the weak one to a flat line -- the frame has to fit the
        strong one. Each trace is scaled so its own peak reaches the same
        fraction of the frame, which is what makes a comparison legible. The
        scales stay adjustable afterwards; this only sets a sane starting
        point.
        """
        visible = self._visible_1d()
        if not visible:
            return
        frame = self._frame_y_limits(visible, 0.0)
        if frame is None:
            return
        span = frame[1] - frame[0]
        if span <= 0:
            return
        # In stacked mode each trace also gets an offset slot, so the share
        # of the frame available to any one of them is a fraction of the whole.
        divisor = len(visible) if self._arrangement == self.ARRANGEMENT_STACKED else 1
        target = (span * target_fraction) / divisor
        for trace in visible:
            peak = float(np.nanmax(np.abs(trace.intensity)))
            if peak > 0 and np.isfinite(peak):
                trace.y_scale = max(MIN_Y_SCALE, min(target / peak, MAX_Y_SCALE))
        self._redraw()
        # Then pin the frame around what is now actually drawn, so every
        # spectrum sits inside the canvas rather than running off the top.
        self.fit_to_drawn()
        self.tracesChanged.emit()

    def fit_to_drawn(self) -> None:
        """Frame exactly what is currently plotted, including scale and offset.

        Distinct from reset_y_limits(), which returns to the neutral raw-data
        frame. This one guarantees everything visible is inside the canvas --
        what "fit all spectra" means to a reader.

        In stacked mode it also re-establishes the lane grid. Scaling a
        spectrum deliberately leaves the grid alone (so its neighbours do not
        move and it clips instead), which means the lanes eventually need
        re-laying to suit the new heights -- and Fit Y is the button a reader
        reaches for when the view needs sorting out. Re-laying and then
        fitting has to happen in that order, or the frame is fitted to
        positions the redraw is about to change.
        """
        # Fit means fit: an explicit Y zoom is exactly what the user is
        # asking to be replaced, so it goes before anything is measured. Left
        # in place, the zoom would be re-applied by the redraw below and the
        # fitted limits discarded -- the button would appear to do nothing.
        self._y_range = None
        if (self._arrangement == self.ARRANGEMENT_STACKED
                and self._stack_step is not None):
            self._stack_step = None
            self._y_limits = None
            self._redraw()

        window = self._visible_ppm_window()
        lows, highs = [], []
        for line in self._axes.lines:
            data = line.get_ydata()
            if data is None or len(data) == 0:
                continue
            arr = np.asarray(data, dtype=float)
            if window is not None:
                # Restricted to the visible ppm range for the same reason as
                # everything else here: fitting to a peak that is off-screen
                # flattens the part being looked at.
                x = np.asarray(line.get_xdata(), dtype=float)
                if x.size == arr.size:
                    inside = (x >= window[0]) & (x <= window[1])
                    if np.any(inside):
                        arr = arr[inside]
            if not np.any(np.isfinite(arr)):
                continue
            lows.append(float(np.nanmin(arr)))
            highs.append(float(np.nanmax(arr)))
        if not lows:
            return
        low, high = min(lows), max(highs)
        if not np.isfinite(low) or not np.isfinite(high):
            return
        pad = (high - low) * 0.06 if high > low else (abs(high) * 0.1 or 1.0)
        self._y_limits = (low - pad, high + pad)
        self._axes.set_ylim(*self._y_limits)
        self._canvas.draw_idle()

    def _visible_ppm_window(self):
        """(low, high) ppm currently on screen, or None for "all of it"."""
        if self._ppm_range:
            left, right = self._ppm_range
            return (min(left, right), max(left, right))
        return None

    def _window_values(self, trace):
        """This trace's intensities WITHIN the visible ppm window.

        Every vertical calculation has to go through here. A 19F spectrum runs
        from -29 to -180 ppm and is routinely examined over a 10 ppm slice; a
        strong peak a hundred ppm outside the view was still setting the frame
        height, the stack lane height and the fit, so the data actually on
        screen was compressed into a flat line -- measured at 0.9% of the
        canvas. Vertical scale must be decided by what is being looked at.

        Falls back to the full array when nothing is in the window, so a
        spectrum scrolled out of view degrades to its old behaviour rather
        than to an exception.
        """
        data = np.asarray(trace.intensity, dtype=float)
        if data.size == 0:
            return data
        window = self._visible_ppm_window()
        if window is None:
            return data
        ppm = np.asarray(trace.ppm, dtype=float)
        if ppm.size != data.size:
            return data
        inside = (ppm >= window[0]) & (ppm <= window[1])
        return data[inside] if np.any(inside) else data

    def _displayed_floor(self, trace) -> float | None:
        """The lowest point of this trace AS DRAWN, within the visible window.

        Two corrections over taking nanmin of the whole array:

        * **scale is applied**, because that is what is on screen. Without it,
          a trace scaled up by "Same noise" (which readily produces factors of
          20 or more) was positioned by its unscaled minimum and ended up far
          off the bottom of the canvas -- reported as "I cannot bottom other
          spectra, nothing shows".
        * **only the visible ppm window counts.** A 19F spectrum displayed
          from -120 to -130 ppm may span -29 to -180 ppm in full, and the
          global minimum can be an artefact hundreds of ppm away from anything
          being looked at. Aligning to a baseline that is not on screen puts
          the visible part somewhere arbitrary.
        """
        data = self._window_values(trace)
        if data.size == 0:
            return None
        floor = float(np.nanmin(data))
        if not np.isfinite(floor):
            return None
        # The gain is applied in stacked mode, so the drawn floor includes
        # it. Without this, "to bottom" positioned a zoomed stack by its
        # unzoomed minimum and put the baseline in the wrong place.
        if self._arrangement == self.ARRANGEMENT_STACKED:
            return floor * trace.y_scale * self._stack_gain
        return floor * trace.y_scale

    def _bottom_anchor(self) -> float | None:
        """The y value a bottomed baseline is placed on.

        Taken from the RAW-data frame, not from the current axis limits.
        Anchoring to the axis meant every press moved the trace to wherever
        the last press had left the frame -- so "To bottom" crept downwards
        each time, and "Bottom All" gave each spectrum a different offset
        because the frame shifted underneath the loop. The raw frame does not
        move when offsets change, which is exactly what makes this repeatable.
        """
        frame = self._frame_y_limits(self._visible_1d(), 0.0)
        return None if frame is None else frame[0]

    def move_to_bottom(self, index: int) -> None:
        """Sit one spectrum's baseline on the bottom of the frame."""
        if not (0 <= index < len(self._traces)):
            return
        trace = self._traces[index]
        if trace.is_2d:
            return   # a contour panel has no baseline to move
        anchor = self._bottom_anchor()
        floor = self._displayed_floor(trace)
        if anchor is None or floor is None:
            return
        self.push_undo()
        trace.y_offset = anchor - floor
        # In stacked mode the drawing adds the stack step for this trace's
        # position, so this lands it on the bottom of ITS OWN lane -- the
        # canvas floor for the first spectrum, halfway up for the second of
        # two, a third of the way for the second of three. No special case is
        # needed here, and adding one would double-count the step.
        if self._arrangement == self.ARRANGEMENT_STACKED:
            self._y_limits = None
            self._stack_step = None
        self._redraw()
        # Then frame what is actually drawn. A scaled-up trace extends far
        # beyond the raw envelope, so without this it lands correctly and is
        # still invisible.
        self.fit_to_drawn()
        self.tracesChanged.emit()

    def move_all_to_bottom(self) -> None:
        """Every spectrum on a common baseline -- what 'overlay' means once
        offsets have been used.

        The anchor is computed ONCE and applied to all of them. Calling
        move_to_bottom in a loop redrew and refitted between traces, so each
        one was aligned to a frame the previous one had just changed: the
        offsets came out different and the baselines did not line up.
        """
        anchor = self._bottom_anchor()
        if anchor is None:
            return
        self.push_undo()
        moved = False
        for trace in self._traces:
            if trace.is_2d:
                continue
            floor = self._displayed_floor(trace)
            if floor is None:
                continue
            trace.y_offset = anchor - floor
            moved = True
        if not moved:
            return
        if self._arrangement == self.ARRANGEMENT_STACKED:
            self._y_limits = None
            self._stack_step = None
        self._redraw()
        self.fit_to_drawn()
        self.tracesChanged.emit()

    # -- right-click export ---------------------------------------------

    # Format -> (file filter, is_vector). Vector formats keep text editable.
    EXPORT_FILTERS = [
        ("PNG image (*.png)", "png"),
        ("JPEG image (*.jpg)", "jpg"),
        ("TIFF image (*.tif)", "tif"),
        ("SVG vector (*.svg)", "svg"),
        ("PDF vector (*.pdf)", "pdf"),
        ("EPS vector (*.eps)", "eps"),
        ("PostScript (*.ps)", "ps"),
        ("PowerPoint slide (*.pptx)", "pptx"),
    ]
    # EMF is deliberately NOT offered. matplotlib cannot write it, and the
    # only route (converting via an external Inkscape install) silently
    # produced nothing when Inkscape was absent, which is worse than not
    # offering the format. SVG, PDF and EPS are vector and universally
    # readable; use one of those and convert downstream if EMF is required.

    def build_context_menu(self):
        """Menu construction only -- never calls exec().

        Split out so tests can trigger the actions without entering a real
        modal loop, which blocks forever in a test run.
        """
        menu = QMenu(self)
        menu.addAction("Save image\u2026", self.request_save_image)
        menu.addSeparator()
        menu.addAction("Auto scale Y", self.autoscale_traces)
        menu.addAction("Fit Y to data", self.fit_to_drawn)
        menu.addAction("All to bottom", self.move_all_to_bottom)
        return menu

    def _on_context_menu(self, pos) -> None:
        if not self._traces:
            return
        self.build_context_menu().exec(self.mapToGlobal(pos))

    def request_save_image(self) -> None:
        """Ask for a path and write the figure exactly as displayed."""
        filters = ";;".join(f for f, _ in self.EXPORT_FILTERS)
        path, chosen = QFileDialog.getSaveFileName(
            self, "Save spectrum image", "", filters
        )
        if not path:
            return
        # If the typed name has no extension, take it from the chosen filter
        # rather than silently writing a PNG named "figure".
        #
        # Path().suffix, not a manual rsplit on "/": Qt returns native
        # separators, so on Windows the whole string came back as the
        # "filename" and any dot anywhere in the path -- "OneDrive - Company
        # Corp", a versioned folder, a dotted user name -- read as an
        # extension already being present, and the figure was written with no
        # extension at all.
        if not Path(path).suffix:
            ext = next(
                (e for f, e in self.EXPORT_FILTERS if f == chosen), "png"
            )
            path = f"{path}.{ext}"
        try:
            self.save_image(path)
        except Exception as exc:  # noqa: BLE001 -- report, never crash
            self.loadFailed.emit(path, f"could not save: {exc}")
            return
        self.imageSaved.emit(path)

    def save_image(self, path, transparent: bool = False, dpi: int = 300) -> None:
        """Save exactly what is on screen.

        The tight-bounding-box save option is deliberately NOT used, and the
        layout is not altered: those silently re-crop and re-scale the figure,
        so the file would not match the canvas. The figure is written at its
        current size and limits, which is what "what you see is what you get"
        has to mean. (The banned keyword is enforced by test_architecture.py.)

        Vector formats (svg, pdf, eps, ps) keep text as editable text rather
        than outlines, so labels can be adjusted downstream.
        """
        import matplotlib as mpl

        # The crosshair is a UI overlay, not part of the figure. Leaving it in
        # bakes a stray cursor line into every exported image.
        self._clear_crosshair()

        path = str(path)
        # Path().suffix, NOT a manual rsplit on ".": a dot anywhere in the
        # path -- a versioned folder, "OneDrive - Company Corp", a dotted user
        # name -- makes rsplit return a fragment of a DIRECTORY name as the
        # format ("2/figure" for /data/v1.2/figure), which matplotlib then
        # rejects. request_save_image already learned this; this is the same
        # bug in its sibling, reachable by any direct caller.
        suffix = Path(path).suffix.lstrip(".").lower() or "png"
        facecolor = "none" if transparent else "white"

        # matplotlib converts text to outlines by default in vector output,
        # which makes labels uneditable in Illustrator/Inkscape -- the whole
        # point of exporting vector. 'none' for SVG keeps real <text>; type 42
        # (TrueType) for PDF/PS/EPS embeds editable text rather than Type 3.
        overrides = {
            "svg.fonttype": "none",
            "pdf.fonttype": 42,
            "ps.fonttype": 42,
        }
        if suffix == "pptx":
            self._save_pptx(path, dpi=dpi)
            return
        with mpl.rc_context(overrides):
            self._figure.savefig(
                path,
                format=suffix,
                dpi=dpi,
                facecolor=facecolor,
                edgecolor="none",
                transparent=transparent,
            )

    def _save_pptx(self, path: str, dpi: int = 300) -> None:
        """One slide sized to the figure, with the plot placed on it.

        The image is embedded as EMF-free PNG at full resolution: PowerPoint
        cannot render SVG reliably across versions, so a high-DPI raster is
        the format that actually survives being opened on someone else's
        machine.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:   # pragma: no cover - dependency present
            raise RuntimeError(
                "PowerPoint export needs the python-pptx package"
            ) from exc

        import tempfile as _tempfile

        width_in, height_in = self._figure.get_size_inches()
        with _tempfile.TemporaryDirectory() as tmp:
            png = f"{tmp}/slide.png"
            self._figure.savefig(
                png, format="png", dpi=dpi, facecolor="white", edgecolor="none"
            )
            prs = Presentation()
            prs.slide_width = Inches(float(width_in))
            prs.slide_height = Inches(float(height_in))
            # 6 = the blank layout in the default template.
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                png, 0, 0,
                width=Inches(float(width_in)),
                height=Inches(float(height_in)),
            )
            prs.save(path)

    def add_spectra(self, index_a: int, index_b: int, label: str | None = None) -> bool:
        """Add a sum spectrum (A + B) as a new trace. Same rules as subtract."""
        return self._combine(index_a, index_b, "+", label)

    def subtract(self, index_a: int, index_b: int, label: str | None = None) -> bool:
        """Add a difference spectrum (A - B) as a new trace.

        The two spectra are INTERPOLATED onto a common ppm axis before
        subtracting. Subtracting index-by-index is the classic silent error
        here: two spectra rarely have identical point counts or sweep widths,
        and index subtraction then quietly compares different chemical shifts
        and produces convincing nonsense.

        Their y_scale is applied first, so a difference taken after using
        "Same noise" or manual scaling reflects what is actually on screen.
        Returns False if the spectra do not overlap in ppm at all.
        """
        return self._combine(index_a, index_b, "-", label)

    def _combine(self, index_a: int, index_b: int, op: str,
                 label: str | None = None) -> bool:
        n = len(self._traces)
        if not (0 <= index_a < n and 0 <= index_b < n) or index_a == index_b:
            return False
        a, b = self._traces[index_a], self._traces[index_b]
        if a.is_2d or b.is_2d:
            return False   # 2D differences are not supported yet

        # y_scale is applied first, so a difference taken after "Same noise"
        # or manual scaling reflects what is actually on screen.
        result = combine_arrays(
            a.ppm, a.intensity, a.y_scale, b.ppm, b.intensity, b.y_scale, op
        )
        if result is None:
            return False   # no overlap in ppm at all
        common, combined = result
        self.push_undo()

        symbol = "+" if op == "+" else "\u2212"
        slot = self._slot_styles[len(self._traces) % len(self._slot_styles)]
        self._traces.append(
            Trace(
                path=Path(f"{a.path}::{op}::{b.path}"),
                label=label or f"\u0394  {a.label}  {symbol}  {b.label}",
                ppm=common,
                intensity=combined,
                color=slot["color"],
                line_width=slot["width"],
                line_style=slot["style"],
                # Marked as derived so the legend, and the session, can tell a
                # difference apart from a measured spectrum.
                is_difference=True,
                source_a=a.path,
                source_b=b.path,
                operation=op,
                source_scales=(a.y_scale, b.y_scale),
            )
        )
        self._y_limits = None
        self._redraw()
        self.tracesChanged.emit()
        return True

    def set_show_pulse_program(self, show: bool) -> None:
        self._show_pulprog = bool(show)
        self._redraw()

    def show_pulse_program(self) -> bool:
        return self._show_pulprog

    def _emit_mode_if_changed(self) -> None:
        current = self.mode()
        if current != self._last_mode:
            self._last_mode = current
            self.modeChanged.emit(current)

    def mode(self) -> str:
        """'2D' when any visible spectrum is 2D, otherwise '1D'.

        Derived from the data rather than being a separate switch the user has
        to remember to flip -- the app already knows each dataset's
        dimensionality.
        """
        return "2D" if any(t.visible and t.is_2d for t in self._traces) else "1D"

    def set_f1_range(self, high: float, low: float) -> None:
        """Indirect-dimension (F1) range, 2D only."""
        if high == low:
            return
        high, low = max(high, low), min(high, low)
        self._f1_range = (high, low)
        self._redraw()

    def set_f2_range(self, high: float, low: float) -> None:
        """Direct-dimension (F2) range, 2D only."""
        if high == low:
            return
        high, low = max(high, low), min(high, low)
        self._f2_range = (high, low)
        self._redraw()

    def f1_bounds(self):
        panels = [t for t in self._traces if t.visible and t.is_2d]
        if not panels:
            return None
        return (
            max(float(np.nanmax(t.ppm_f1)) for t in panels),
            min(float(np.nanmin(t.ppm_f1)) for t in panels),
        )

    def f2_bounds(self):
        panels = [t for t in self._traces if t.visible and t.is_2d]
        if not panels:
            return None
        return (
            max(float(np.nanmax(t.ppm_f2)) for t in panels),
            min(float(np.nanmin(t.ppm_f2)) for t in panels),
        )

    def display_label(self, trace) -> str:
        """Name shown on the plot: sample/expno plus the pulse programme.

        The pulse programme is what distinguishes two experiments on the same
        sample, which is exactly the case where a bare name is ambiguous.
        """
        if trace.pulse_program:
            return f"{trace.label}  ({trace.pulse_program})"
        return trace.label

    def set_label_offset(self, index: int, dx: float, dy: float) -> None:
        """Move one spectrum's name by an explicit amount, in axes fractions.

        Dragging is available too, but a typed offset is reproducible and
        works even when a name sits underneath a trace and is awkward to
        grab.
        """
        if not (0 <= index < len(self._traces)):
            return
        try:
            dx = float(dx)
            dy = float(dy)
        except (TypeError, ValueError):
            return
        if not (np.isfinite(dx) and np.isfinite(dy)):
            return
        trace = self._traces[index]
        base = trace.label_base_pos or trace.label_pos or (0.015, 0.985)
        trace.label_base_pos = base
        trace.label_offset = (dx, dy)
        trace.label_pos = (
            min(max(base[0] + dx, 0.0), 0.98),
            min(max(base[1] + dy, 0.02), 1.0),
        )
        self._redraw()
        self.tracesChanged.emit()

    def reset_label_positions(self) -> None:
        """Put dragged spectrum names back in the default stacked column."""
        for trace in self._traces:
            trace.label_pos = None
            trace.label_offset = (0.0, 0.0)
            trace.label_base_pos = None
            trace.label_dx = 0.0
            trace.label_dy = 0.0
        self._redraw()

    def set_label_scale(self, scale: float) -> None:
        """Size of the on-plot spectrum names, relative to default."""
        try:
            scale = float(scale)
        except (TypeError, ValueError):
            return
        if not np.isfinite(scale) or not (0.3 <= scale <= 4.0):
            return
        self._label_scale = scale
        self._redraw()

    def label_scale(self) -> float:
        return self._label_scale

    def set_x_decimals(self, decimals) -> None:
        """Decimal places on the ppm axis (0 -> "1", 2 -> "1.00"). None =
        let matplotlib choose."""
        if decimals is not None:
            try:
                decimals = int(decimals)
            except (TypeError, ValueError):
                return
            if decimals < 0 or decimals > 6:
                return
        self._x_decimals = decimals
        self._redraw()

    def x_decimals(self):
        return self._x_decimals

    def normalise_to_noise(self, target_snr_reference: int = 0) -> bool:
        """Scale every spectrum so their NOISE levels match.

        Equalising noise is what makes peak heights comparable between spectra
        acquired with different numbers of scans or receiver gain -- a peak
        twice as tall then genuinely means twice the signal-to-noise, which is
        the comparison a reader actually wants to make.

        Noise is estimated robustly as the median absolute deviation of the
        whole trace. MAD is used rather than the standard deviation because
        the peaks themselves are outliers: a plain sigma is dominated by the
        signal and would be a measure of peak height, not noise.

        Returns False when the noise cannot be estimated (e.g. a perfectly
        flat trace), rather than silently scaling by a meaningless factor.
        """
        visible = self._visible_1d()
        if len(visible) < 1:
            return False
        if not (0 <= target_snr_reference < len(visible)):
            target_snr_reference = 0

        def noise_of(trace) -> float:
            data = np.asarray(trace.intensity, dtype=float)
            if data.size == 0:
                return 0.0
            median = float(np.nanmedian(data))
            mad = float(np.nanmedian(np.abs(data - median)))
            # 1.4826 converts MAD to a Gaussian-equivalent sigma.
            return mad * 1.4826

        noises = [noise_of(t) for t in visible]
        reference = noises[target_snr_reference]
        if reference <= 0 or not np.isfinite(reference):
            return False

        for trace, noise in zip(visible, noises):
            if noise <= 0 or not np.isfinite(noise):
                continue
            factor = reference / noise
            trace.y_scale = max(MIN_Y_SCALE, min(factor, MAX_Y_SCALE))
        self._redraw()
        self.fit_to_drawn()
        self.tracesChanged.emit()
        return True

    # -- session state ------------------------------------------------------

    SESSION_FORMAT = 1

    def session_state(self) -> dict:
        """Everything needed to rebuild this view later.

        Spectra are stored as PATHS, not as their data: the arrays are large,
        they already exist on disk, and re-reading them keeps a saved session
        honest if the underlying processing is redone. A difference has no
        file behind it, so what is stored is the RECIPE -- which two spectra,
        which operator, and the scales they carried when it was made -- and
        restore re-derives it. Same principle, one level up. Everything the user
        adjusted -- scale, offset, colour, style, label position, ranges,
        contour settings -- is stored, because that is the part that cannot be
        recovered by re-reading.
        """
        return {
            "format": self.SESSION_FORMAT,
            "arrangement": self._arrangement,
            "ppm_range": list(self._ppm_range) if self._ppm_range else None,
            "y_range": list(self._y_range) if self._y_range else None,
            "stack_gain": self._stack_gain,
            "f1_range": list(self._f1_range) if self._f1_range else None,
            "f2_range": list(self._f2_range) if self._f2_range else None,
            "grid": self._show_grid,
            "grid_spacing_ppm": self._grid_spacing_ppm,
            "grid_spacing_y": self._grid_spacing_y,
            "x_decimals": self._x_decimals,
            "label_scale": self._label_scale,
            "slot_styles": self.slot_styles(),
            "selected_index": self._selected_index,
            "spectra": [
                {
                    "path": str(t.path),
                    "label": t.label,
                    "is_2d": t.is_2d,
                    "is_difference": t.is_difference,
                    "source_a": str(t.source_a) if t.source_a else None,
                    "source_b": str(t.source_b) if t.source_b else None,
                    "operation": t.operation,
                    "source_scales": list(t.source_scales),
                    "pulse_program": t.pulse_program,
                    "nucleus": t.nucleus,
                    "ns": t.ns,
                    "rg": t.rg,
                    "color": t.color,
                    "line_width": t.line_width,
                    "line_style": t.line_style,
                    "visible": t.visible,
                    "y_scale": t.y_scale,
                    "y_offset": t.y_offset,
                    "x_offset": t.x_offset,
                    "label_pos": list(t.label_pos) if t.label_pos else None,
                    "label_offset": list(t.label_offset),
                    "contour_levels": t.contour_levels,
                    "contour_factor": t.contour_factor,
                    "contour_base_sigma": t.contour_base_sigma,
                }
                for t in self._traces
            ],
        }

    def restore_session(self, state: dict) -> list:
        """Rebuild a saved view. Returns the labels that could not be reloaded.

        Derived spectra (differences and sums) are RE-DERIVED from their
        recorded recipe, in place, so a session round-trips exactly what was
        on screen. They used to be skipped, which silently threw away the
        result of every Subtract the moment a session was saved.

        Order is preserved by filling a slot per saved entry and only then
        compacting: a difference can sit anywhere in the list (Bottom All
        reorders traces), so its sources are not necessarily above it.

        Missing files are reported rather than silently dropped -- a session
        that quietly loses half its spectra is worse than one that says so.
        """
        if not isinstance(state, dict):
            raise ValueError("session file is not a HelSpin session")
        if int(state.get("format", 0)) != self.SESSION_FORMAT:
            raise ValueError(
                f"unsupported session format {state.get('format')!r}"
            )

        self.clear()
        failed = []
        entries = list(state.get("spectra", []))
        slots: list = [None] * len(entries)
        derived: list = []
        for position, entry in enumerate(entries):
            if entry.get("is_difference"):
                derived.append((position, entry))
                continue
            path = Path(entry["path"])
            try:
                if entry.get("is_2d"):
                    spec = self._reader.read_2d(path)
                    trace = Trace(
                        path=path, label=entry["label"],
                        ppm=np.asarray([]), intensity=np.asarray([]),
                        color=entry["color"], is_2d=True,
                        matrix=np.asarray(spec.real, dtype=np.float64),
                        ppm_f1=np.asarray(spec.axis_f1.ppm_scale(), dtype=np.float64),
                        ppm_f2=np.asarray(spec.axis_f2.ppm_scale(), dtype=np.float64),
                    )
                else:
                    spec = self._reader.read_1d(path)
                    trace = Trace(
                        path=path, label=entry["label"],
                        ppm=np.asarray(spec.axis.ppm_scale(), dtype=np.float64),
                        intensity=np.asarray(spec.real, dtype=np.float64),
                        color=entry["color"],
                    )
            except Exception as exc:   # noqa: BLE001 - report, never abort
                failed.append(f"{entry.get('label', path)}: {exc}")
                continue

            trace.line_width = float(entry.get("line_width", 0.8))
            trace.line_style = entry.get("line_style", "-")
            trace.visible = bool(entry.get("visible", True))
            trace.y_scale = float(entry.get("y_scale", 1.0))
            trace.y_offset = float(entry.get("y_offset", 0.0))
            # Absent from sessions written before this feature existed; 0.0 is
            # the correct reading of an older file, which carried no shift.
            trace.x_offset = float(entry.get("x_offset", 0.0))
            pos = entry.get("label_pos")
            trace.label_pos = tuple(pos) if pos else None
            trace.label_offset = tuple(entry.get("label_offset", (0.0, 0.0)))
            trace.pulse_program = entry.get("pulse_program", "")
            trace.nucleus = entry.get("nucleus", "")
            trace.ns = int(entry.get("ns", 0) or 0)
            trace.rg = float(entry.get("rg", 0.0) or 0.0)
            trace.contour_levels = int(entry.get("contour_levels", 12))
            trace.contour_factor = float(entry.get("contour_factor", 1.3))
            trace.contour_base_sigma = float(entry.get("contour_base_sigma", 4.0))
            slots[position] = trace

        for position, entry in derived:
            trace = self._rederive(entry, slots)
            if trace is None:
                failed.append(
                    f"{entry.get('label', 'difference')}: "
                    "its source spectra are not in this session"
                )
                continue
            slots[position] = trace

        self._traces = [t for t in slots if t is not None]

        # Opening a session starts a new document, so its history starts
        # empty. Leaving the previous stack in place meant a second Ctrl+Z
        # after opening a file swapped the whole restored session out for
        # whatever had been on the canvas beforehand.
        self._undo.clear()
        self._redo.clear()
        self._last_undo_key = None
        self.historyChanged.emit()

        styles = state.get("slot_styles")
        if styles:
            self._slot_styles = list(styles)
        self._arrangement = state.get("arrangement", self.ARRANGEMENT_OVERLAY)
        rng = state.get("ppm_range")
        self._ppm_range = tuple(rng) if rng else None
        # Absent in sessions written before Y zoom existed. Absent means
        # automatic framing, which is what those files actually had.
        yr = state.get("y_range")
        self._y_range = tuple(yr) if yr else None
        # Absent in sessions written before stacked Y zoom existed; 1.0
        # is the correct reading of a file that carried no gain.
        try:
            gain = float(state.get("stack_gain", 1.0))
        except (TypeError, ValueError):
            gain = 1.0
        if not np.isfinite(gain) or gain <= 0:
            gain = 1.0
        # Accepted at whatever magnification it was saved at -- see the wheel
        # handler: there is no ceiling. Already rejected above if it is not a
        # finite positive number.
        self._stack_gain = gain
        f1 = state.get("f1_range")
        self._f1_range = tuple(f1) if f1 else None
        f2 = state.get("f2_range")
        self._f2_range = tuple(f2) if f2 else None
        self._show_grid = bool(state.get("grid", False))
        self._grid_spacing_ppm = state.get("grid_spacing_ppm")
        self._grid_spacing_y = state.get("grid_spacing_y")
        self._x_decimals = state.get("x_decimals")
        self._label_scale = float(state.get("label_scale", 1.0))
        index = state.get("selected_index")
        self._selected_index = (
            index if isinstance(index, int) and 0 <= index < len(self._traces)
            else (0 if self._traces else None)
        )

        # Y limits are intentionally NOT restored: they are derived from the
        # data, and the saved pair could be stale if the files were reprocessed.
        self._y_limits = None
        self._redraw()
        self._emit_mode_if_changed()
        self.tracesChanged.emit()
        return failed

    def _rederive(self, entry: dict, slots: list):
        """Rebuild one saved difference from its sources. None if impossible.

        The sources are matched by path against the traces this session has
        already restored, so a difference survives only alongside the spectra
        it was made from -- which is the honest outcome: without them there is
        nothing to recompute, and inventing an array would be worse.
        """
        source_a = entry.get("source_a")
        source_b = entry.get("source_b")
        if not source_a or not source_b:
            return None        # saved before 0.4.2: no recipe was recorded
        by_path = {
            str(t.path): t for t in slots if t is not None and not t.is_2d
        }
        a, b = by_path.get(str(source_a)), by_path.get(str(source_b))
        if a is None or b is None:
            return None

        scales = entry.get("source_scales") or [a.y_scale, b.y_scale]
        try:
            scale_a, scale_b = float(scales[0]), float(scales[1])
        except (TypeError, ValueError, IndexError):
            scale_a, scale_b = a.y_scale, b.y_scale

        result = combine_arrays(
            a.ppm, a.intensity, scale_a,
            b.ppm, b.intensity, scale_b,
            entry.get("operation", "-"),
        )
        if result is None:
            return None
        common, combined = result

        trace = Trace(
            path=Path(entry["path"]),
            label=entry["label"],
            ppm=common,
            intensity=combined,
            color=entry["color"],
            is_difference=True,
            source_a=Path(source_a),
            source_b=Path(source_b),
            operation=entry.get("operation", "-"),
            source_scales=(scale_a, scale_b),
        )
        trace.line_width = float(entry.get("line_width", 0.8))
        trace.line_style = entry.get("line_style", "-")
        trace.visible = bool(entry.get("visible", True))
        trace.y_scale = float(entry.get("y_scale", 1.0))
        trace.y_offset = float(entry.get("y_offset", 0.0))
        trace.x_offset = float(entry.get("x_offset", 0.0))
        pos = entry.get("label_pos")
        trace.label_pos = tuple(pos) if pos else None
        trace.label_offset = tuple(entry.get("label_offset", (0.0, 0.0)))
        return trace

    def set_contour_levels(self, index: int, levels: int) -> None:
        """How many contour lines to draw for one 2D spectrum."""
        if not (0 <= index < len(self._traces)):
            return
        try:
            levels = int(levels)
        except (TypeError, ValueError):
            return
        if not (2 <= levels <= 60):
            return
        self._traces[index].contour_levels = levels
        self._redraw()
        self.tracesChanged.emit()

    def set_contour_factor(self, index: int, factor: float) -> None:
        """Ratio between successive contour levels."""
        if not (0 <= index < len(self._traces)):
            return
        try:
            factor = float(factor)
        except (TypeError, ValueError):
            return
        if not np.isfinite(factor) or not (1.01 <= factor <= 5.0):
            return
        self._traces[index].contour_factor = factor
        self._redraw()
        self.tracesChanged.emit()

    def set_contour_base_sigma(self, index: int, sigma: float) -> None:
        """Lowest contour, as a multiple of the noise estimate."""
        if not (0 <= index < len(self._traces)):
            return
        try:
            sigma = float(sigma)
        except (TypeError, ValueError):
            return
        if not np.isfinite(sigma) or not (0.5 <= sigma <= 50.0):
            return
        self._traces[index].contour_base_sigma = sigma
        self._redraw()
        self.tracesChanged.emit()

    def apply_contour_defaults(self, levels: int, factor: float,
                               base_sigma: float) -> None:
        """Apply contour settings to every 2D spectrum at once."""
        for i, trace in enumerate(self._traces):
            if trace.is_2d:
                self.set_contour_levels(i, levels)
                self.set_contour_factor(i, factor)
                self.set_contour_base_sigma(i, base_sigma)

    def _apply_grid(self, axes, horizontal: bool = True) -> None:
        """Faint reference grid on one or both axes.

        Deliberately faint: a grid on a spectrum is a reading aid, not a
        feature of the data, so it must not compete with the peaks.

        Shared by the 1D and both 2D paths. The 2D drawing had no grid code at
        all, so the toolbar toggle silently did nothing there -- the button
        was on, and no grid appeared.
        """
        from matplotlib.ticker import MultipleLocator

        if self._grid_spacing_ppm:
            axes.xaxis.set_major_locator(MultipleLocator(self._grid_spacing_ppm))
        axes.grid(
            True, axis="x", which="major",
            linewidth=0.4, alpha=0.25, linestyle="-",
        )
        if horizontal:
            if self._grid_spacing_y:
                axes.yaxis.set_major_locator(
                    MultipleLocator(self._grid_spacing_y)
                )
            axes.grid(
                True, axis="y", which="major",
                linewidth=0.4, alpha=0.25, linestyle="-",
            )
        else:
            axes.grid(False, axis="y")

    def set_grid_spacing_y(self, spacing) -> None:
        """Spacing for the horizontal grid: F1 ppm in 2D, intensity in 1D.

        Separate from the x spacing because in 2D the two axes cover quite
        different ranges -- 10 ppm of proton against 150 of carbon -- so one
        number cannot suit both.
        """
        try:
            spacing = float(spacing)
        except (TypeError, ValueError):
            return
        self._grid_spacing_y = spacing if spacing > 0 else None
        self._redraw()

    def grid_spacing_y(self):
        return self._grid_spacing_y

    def set_grid_spacing_ppm(self, spacing) -> None:
        """Fixed grid spacing in ppm, or None to let matplotlib choose."""
        if spacing is not None:
            try:
                spacing = float(spacing)
            except (TypeError, ValueError):
                return
            if spacing <= 0 or not np.isfinite(spacing):
                return
        self._grid_spacing_ppm = spacing
        self._redraw()

    def grid_spacing_ppm(self):
        return self._grid_spacing_ppm

    def reset_y_limits(self) -> None:
        """Recompute the y range from what is currently drawn.

        Drops any explicit Y zoom as well. Without that this did nothing at
        all once the user had zoomed -- it cleared the cache, the redraw
        re-applied the zoom on top, and the control looked broken.

        Undoable, but only when there is something to lose. Before Y zoom
        existed this cleared a CACHE, and undoing a cache clear is
        meaningless; now it can discard a vertical window the user set
        deliberately, and one stray click on Fit Y must not be the end of it.
        Pushing unconditionally would be the other error: it would fill the
        undo history with steps that restore nothing visible.

        The key is its own, not "y_zoom". Sharing that key would let a Fit Y
        within the coalescing window fold into the wheel burst that preceded
        it, so a single undo would jump back past both and the intermediate
        zoom -- the state the user is actually trying to recover -- would
        never be reachable.
        """
        if self._y_range is not None or self._stack_gain != 1.0:
            self.push_undo("y_fit")
        self._y_range = None
        self._y_limits = None
        # Fit Y is the single escape hatch for the vertical axis, so it has to
        # clear the stacked gain as well. Leaving it set would make the button
        # appear to do nothing in a stacked view that had been zoomed -- the
        # same fault this method was already fixed for once, with _y_range.
        self._stack_gain = 1.0
        self._redraw()

    def apply_styles(self, styles: list[dict]) -> None:
        """Install per-slot appearance and re-style existing spectra.

        Slot i applies to the i-th loaded spectrum, so the mapping is the same
        before and after anything is dropped.
        """
        if not styles:
            return
        self._slot_styles = list(styles)
        for i, trace in enumerate(self._traces):
            slot = self._slot_styles[i % len(self._slot_styles)]
            trace.color = slot["color"]
            trace.line_style = slot["style"]
            trace.line_width = slot["width"]
        self._redraw()
        self.tracesChanged.emit()

    def slot_styles(self) -> list[dict]:
        return [dict(s) for s in self._slot_styles]

    def set_grid_visible(self, visible: bool) -> None:
        self._show_grid = bool(visible)
        self._redraw()

    def grid_visible(self) -> bool:
        return self._show_grid

    def set_trace_line_style(self, index: int, style: str) -> None:
        """Solid / dashed / dotted / dash-dot, per trace."""
        if not (0 <= index < len(self._traces)):
            return
        if style not in LINE_STYLES.values():
            return
        self._traces[index].line_style = style
        self._redraw()

    # -- cursor crosshair and drag-to-move -----------------------------------

    def set_crosshair_enabled(self, enabled: bool) -> None:
        self._crosshair_enabled = bool(enabled)
        if not enabled:
            self._clear_crosshair()

    def _clear_crosshair(self) -> None:
        if self._crosshair is not None:
            for artist in self._crosshair:
                try:
                    artist.remove()
                except (ValueError, NotImplementedError):
                    pass
            self._crosshair = None
            self._canvas.draw_idle()

    def _on_mouse_move(self, event) -> None:
        if self._zoom_mode and self._zoom_band is not None:
            self._update_zoom_band(event)
            return
        # Dragging the selected trace takes precedence over the crosshair.
        if self._label_drag is not None:
            self._continue_label_drag(event)
            return
        if self._drag_start is not None and event.inaxes is self._axes:
            self._continue_drag(event)
            return
        if event.inaxes is not self._axes:
            self._clear_crosshair()
            return
        if not getattr(self, "_crosshair_enabled", True):
            return
        self._clear_crosshair()
        vline = self._axes.axvline(
            event.xdata, color="#888888", linewidth=0.6, linestyle="--", alpha=0.8
        )
        hline = self._axes.axhline(
            event.ydata, color="#888888", linewidth=0.6, linestyle="--", alpha=0.8
        )
        # The ppm value at the cursor, drawn on the plot itself -- reading it
        # off the axis by eye is exactly what a crosshair is meant to avoid.
        # No "ppm" on the plot labels: the axis directly beneath already says
        # it, and the crosshair readouts are where space is tightest.
        digits = self._cursor_decimals
        label = self._axes.text(
            event.xdata, 1.005, f"{event.xdata:.{digits}f}",
            # x in DATA coords (follows the cursor), y in AXES coords (pinned
            # just above the frame) -- that is what get_xaxis_transform gives.
            transform=self._axes.get_xaxis_transform(),
            ha="center", va="bottom", fontsize=8, color="#333333",
            clip_on=False,
        )
        # The same thing for the horizontal line, at the right-hand edge. The
        # vertical line had carried its value since the crosshair existed and
        # the horizontal one had not, so in 2D -- where both are chemical
        # shifts and the indirect dimension is often the one being read -- half
        # the crosshair was decoration. In 1D the second number is an
        # intensity, so it is formatted as one rather than labelled ppm.
        if self.mode() == "2D":
            y_text = f"{event.ydata:.{digits}f}"
        else:
            y_text = f"{event.ydata:.4g}"      # an intensity, not a shift
        y_label = self._axes.text(
            1.005, event.ydata, y_text,
            # Mirror of the above: y in DATA coords so it tracks the line,
            # x in AXES coords so it stays pinned outside the right spine.
            transform=self._axes.get_yaxis_transform(),
            ha="left", va="center", fontsize=8, color="#333333",
            clip_on=False,
            # A backing box, so the value stays readable if a long number or a
            # narrow window pushes it back over the frame.
            bbox={"facecolor": "white", "edgecolor": "none",
                  "alpha": 0.75, "pad": 1.0},
        )
        self._crosshair = (vline, hline, label, y_label)
        self.cursorMoved.emit(float(event.xdata), float(event.ydata))
        self._canvas.draw_idle()

    def _on_mouse_press(self, event) -> None:
        """Begin dragging the selected trace vertically.

        Only in stacked mode and only with a trace selected: dragging in
        overlay would be ambiguous about which spectrum is meant, and moving
        one without a clear selection would feel arbitrary.
        """
        if event.button != 1:
            return
        if self._zoom_mode:
            # In zoom mode the drag draws a box; it must not also move a
            # trace, or one gesture would do two things at once.
            if event.inaxes is not None:
                self._begin_zoom_band(event)
            return
        if event.inaxes is not self._axes:
            return
        # Dragging works in BOTH arrangements. Restricting it to stacked (an
        # earlier decision) meant the control silently did nothing in overlay,
        # which reads as broken rather than as a deliberate limitation. The
        # selection makes the target unambiguous either way.
        # A click on a spectrum NAME drags the name, not the spectrum -- the
        # names are what overlap the data and need moving out of the way.
        hit = self._label_at(event)
        if hit is not None:
            artist, hit_trace = hit
            self._label_drag = (hit_trace, event.x, event.y, hit_trace.label_pos)
            return

        trace = self.selected_trace()
        if trace is None or event.ydata is None:
            return
        self._drag_start = (float(event.ydata), trace.y_offset)

    def _label_at(self, event):
        """The spectrum name under the cursor, if any."""
        for artist, trace in getattr(self, "_label_artists", []):
            try:
                contains, _info = artist.contains(event)
            except (TypeError, AttributeError):
                continue
            if contains:
                return artist, trace
        return None

    def _continue_label_drag(self, event) -> None:
        trace, x0, y0, start_pos = self._label_drag
        if event.x is None or event.y is None:
            return
        width = max(self._canvas.width(), 1)
        height = max(self._canvas.height(), 1)
        dx = (event.x - x0) / width
        dy = (event.y - y0) / height
        nx = min(max(start_pos[0] + dx, 0.0), 0.98)
        ny = min(max(start_pos[1] + dy, 0.02), 1.0)
        trace.label_pos = (nx, ny)
        self._redraw()

    def _continue_drag(self, event) -> None:
        if event.ydata is None:
            return
        start_y, start_offset = self._drag_start
        index = self._selected_index
        if index is None:
            return
        self._traces[index].y_offset = start_offset + (float(event.ydata) - start_y)
        self._redraw()

    def _on_mouse_release(self, event) -> None:
        if self._zoom_mode and self._zoom_band is not None:
            self._finish_zoom_band(event)
            return
        if self._label_drag is not None:
            self._label_drag = None
            self.tracesChanged.emit()
            return
        if self._drag_start is not None:
            self._drag_start = None
            self.tracesChanged.emit()

    def _on_mouse_leave(self, event) -> None:
        self._clear_crosshair()

    # -- drag and drop -------------------------------------------------------
    #
    # dragEnterEvent MUST call acceptProposedAction() or dropEvent never fires
    # -- the single most common Qt drag-and-drop mistake.

    def dragEnterEvent(self, event) -> None:
        if event.mimeData().hasFormat(MIME_DATASET):
            event.acceptProposedAction()

    def dragMoveEvent(self, event) -> None:
        if event.mimeData().hasFormat(MIME_DATASET):
            event.acceptProposedAction()

    def dropEvent(self, event) -> None:
        if self.handle_mime_data(event.mimeData()):
            event.acceptProposedAction()

    def handle_mime_data(self, mime) -> bool:
        """Decode a dataset payload and queue every entry.

        Separate from dropEvent so tests can exercise it with a plain
        QMimeData, without constructing QDropEvent objects whose signature
        varies between Qt versions.
        """
        if not mime.hasFormat(MIME_DATASET):
            return False
        try:
            payload = json.loads(bytes(mime.data(MIME_DATASET)).decode("utf-8"))
        except (ValueError, UnicodeDecodeError):
            return False
        if not payload:
            return False
        for item in payload:
            self.add_dataset(
                item["path"],
                item.get("label") or Path(item["path"]).name,
                int(item.get("dimensionality", 1)),
                pulse_program=item.get("pulse_program", ""),
                nucleus=item.get("nucleus", ""),
            )
        return True
